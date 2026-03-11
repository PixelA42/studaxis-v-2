"""
Studaxis FastAPI Backend — API Bridge

Serves as the bridge between the React frontend and the local AI engine (Ollama),
flashcards system, and optional cloud. CORS is configured for local React development.

When frontend/dist exists (e.g. repo root frontend/dist), serves the compiled React SPA
at / with fallback to index.html for client-side routes. API routes are under /api.

Run: uvicorn main:app --reload --host 0.0.0.0 --port 6782
(from the backend directory)

Or from repo root: uvicorn backend.main:app --reload --host 0.0.0.0 --port 6782
Or: python main.py  (root main.py mounts SPA and runs uvicorn)
"""

from __future__ import annotations

import json
import logging
import os
import re
import time
import uuid
from datetime import datetime, timezone
from pathlib import Path
from typing import Annotated, Any, Optional

from fastapi import Body, Depends, FastAPI, File, Form, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel, Field

# Ensure backend is on path for imports when run from repo root
_APP_DIR = Path(__file__).resolve().parent
if str(_APP_DIR) not in __import__("sys").path:
    __import__("sys").path.insert(0, str(_APP_DIR))

# Load .env early so AWS and AppSync config are available
try:
    from dotenv import load_dotenv
    _env = _APP_DIR / ".env"
    if _env.exists():
        load_dotenv(_env)
except ImportError:
    pass

from ai_integration_layer import AIEngine, AIState, AITaskType
from model_config import get_best_model, get_config_path_for_log
from hardware_validator import ensure_ollama_serve, ensure_ollama_model
from grading.grader import Grader
from grading.red_pen_feedback import RedPenFeedback
from auth_routes import router as auth_router
from database import User, init_db
from dependencies import get_current_user, get_user_id
from profile_store import UserProfile, load_profile, save_profile, load_profile_for_user, save_profile_for_user
from recommendation_service import (
    _has_flashcard_topic,
    _has_quiz_data,
    _get_quiz_profile,
    build_flashcard_based_prompt,
    build_quiz_only_prompt,
    parse_ai_response,
)
from stats_algorithms import (
    ensure_flashcard_structure,
    ensure_streak_structure,
    update_flashcard_entry,
    update_flashcard_stats_from_cards,
    update_quiz_stats as _update_quiz_stats,
    update_streak as _update_streak,
)

# ---------------------------------------------------------------------------
# Base path for AI engine and data (user_stats, profile, etc.)
# When run as "uvicorn main:app", cwd is backend; when "uvicorn backend.main:app", cwd is repo root.
BASE_PATH = Path(os.environ.get("STUDAXIS_BASE_PATH", str(_APP_DIR)))
DATA_DIR = BASE_PATH / "data"
STATS_FILE = DATA_DIR / "user_stats.json"
FLASHCARDS_FILE = DATA_DIR / "flashcards.json"
SAMPLE_TEXTBOOKS_DIR = DATA_DIR / "sample_textbooks"

logger = logging.getLogger("studaxis.main")


def _user_dir(user_id: str) -> Path:
    """Return per-user data directory, creating it if needed."""
    d = DATA_DIR / "users" / user_id
    d.mkdir(parents=True, exist_ok=True)
    return d


def _stats_file(user_id: str) -> Path:
    """Return path to per-user user_stats.json."""
    return _user_dir(user_id) / "user_stats.json"


def _flashcards_file(user_id: str) -> Path:
    """Return path to per-user flashcards.json (legacy)."""
    return _user_dir(user_id) / "flashcards.json"


# New deck-based flashcard storage: data/flashcards/{user_id}.json or {user_id}_{class_id}.json
# Partition by class_id so changing class does not mix flashcards from old and new class (ghost data fix).
def _flashcard_decks_file(user_id: str, class_id: str | None = None) -> Path:
    """Return path to per-user flashcard decks. When class_id set, partition by class to avoid ghost data."""
    d = DATA_DIR / "flashcards"
    d.mkdir(parents=True, exist_ok=True)
    if class_id and str(class_id).strip():
        safe = re.sub(r"[^a-zA-Z0-9_-]", "_", class_id.strip())[:64]
        return d / f"{user_id}_{safe}.json"
    return d / f"{user_id}.json"


def _load_flashcard_decks(user_id: str, class_id: str | None = None) -> list[dict[str, Any]]:
    """Load decks from new structure. Partitioned by class_id when set (ghost data fix)."""
    if class_id is None:
        p = load_profile_for_user(user_id)
        class_id = getattr(p, "class_id", None) if p else None
    decks_path = _flashcard_decks_file(user_id, class_id)
    legacy_path = _flashcards_file(user_id)

    # Migrate legacy flat cards to deck format if new file missing
    if not decks_path.exists() and legacy_path.exists():
        try:
            raw = legacy_path.read_text(encoding="utf-8")
            legacy = json.loads(raw)
            cards = legacy if isinstance(legacy, list) else []
            if cards:
                by_topic: dict[str, list] = {}
                for c in cards:
                    topic = str(c.get("topic") or "General")
                    by_topic.setdefault(topic, []).append(c)
                decks = []
                for topic, topic_cards in by_topic.items():
                    deck_id = f"deck_{uuid.uuid4().hex[:12]}"
                    created = datetime.now(timezone.utc).strftime("%Y-%m-%d")
                    decks.append({
                        "id": deck_id,
                        "title": topic,
                        "subject": topic,
                        "created_at": created,
                        "cards": [_normalize_card_for_deck(c) for c in topic_cards],
                    })
                data = {"decks": decks}
                decks_path.write_text(
                    json.dumps(data, indent=2, ensure_ascii=False),
                    encoding="utf-8",
                )
        except (OSError, json.JSONDecodeError):
            pass

    if not decks_path.exists():
        return []

    try:
        raw = decks_path.read_text(encoding="utf-8")
        data = json.loads(raw)
        return data.get("decks", []) if isinstance(data, dict) else []
    except (OSError, json.JSONDecodeError):
        return []


def _normalize_card_for_deck(c: dict[str, Any]) -> dict[str, Any]:
    """Ensure card has required deck fields; preserve extra fields."""
    base = {
        "id": str(c.get("id") or uuid.uuid4().hex),
        "front": str(c.get("front") or c.get("question") or ""),
        "back": str(c.get("back") or c.get("answer") or ""),
        "ease": str(c.get("ease") or "medium"),
        "next_review": str(c.get("next_review") or datetime.now(timezone.utc).strftime("%Y-%m-%d")),
        "review_count": int(c.get("review_count", 0)),
        "mastered": bool(c.get("mastered", False)),
    }
    for k, v in c.items():
        if k not in base and v is not None:
            base[k] = v
    return base


def _save_flashcard_decks(decks: list[dict[str, Any]], user_id: str, class_id: str | None = None) -> None:
    """Save decks to new structure. Partitioned by class_id when set (ghost data fix)."""
    if class_id is None:
        p = load_profile_for_user(user_id)
        class_id = getattr(p, "class_id", None) if p else None
    try:
        path = _flashcard_decks_file(user_id, class_id)
        path.write_text(
            json.dumps({"decks": decks}, indent=2, ensure_ascii=False),
            encoding="utf-8",
        )
    except OSError:
        pass


def _enqueue_sync(base_path: Path, user_id: str, sync_type: str, payload: dict[str, Any]) -> None:
    """Add item to sync queue for AWS when online."""
    try:
        from sync_manager import SyncManager
        sm = SyncManager(base_path=str(base_path), user_id=user_id)
        if sync_type == "flashcard_review":
            sm._enqueue_generic("flashcard_review", payload)
        elif sync_type == "quiz_result":
            sm._enqueue_generic("quiz_result", payload)
        elif sync_type == "flashcard_create":
            sm._enqueue_generic("flashcard_create", payload)
        elif sync_type == "assignment_complete":
            sm._enqueue_generic("assignment_complete", payload)
    except Exception:
        pass


def _chat_history_file(user_id: str) -> Path:
    """Return path to per-user chat_history.json."""
    return _user_dir(user_id) / "chat_history.json"


def _notifications_file(user_id: str) -> Path:
    """Return path to per-user notifications.json."""
    d = DATA_DIR / "notifications"
    d.mkdir(parents=True, exist_ok=True)
    return d / f"{user_id}.json"


def _assignments_dir() -> Path:
    d = DATA_DIR / "assignments"
    d.mkdir(parents=True, exist_ok=True)
    return d


def _assignments_file(class_code: str) -> Path:
    safe = re.sub(r"[^a-zA-Z0-9_-]", "_", class_code or "default")[:64]
    return _assignments_dir() / f"{safe}.json"


def _load_assignments(class_code: str) -> list[dict[str, Any]]:
    try:
        path = _assignments_file(class_code)
        if path.is_file():
            data = json.loads(path.read_text(encoding="utf-8"))
            return data if isinstance(data, list) else []
    except (OSError, json.JSONDecodeError):
        pass
    return []


def _save_assignments(class_code: str, items: list[dict[str, Any]]) -> None:
    try:
        _assignments_file(class_code).write_text(
            json.dumps(items, indent=2, ensure_ascii=False),
            encoding="utf-8",
        )
    except OSError:
        pass


def _load_notifications(user_id: str) -> list[dict[str, Any]]:
    """Load notifications from JSON file; return [] on missing/error."""
    try:
        f = _notifications_file(user_id)
        if f.exists():
            raw = f.read_text(encoding="utf-8")
            data = json.loads(raw)
            if isinstance(data, list):
                return data
            if isinstance(data, dict) and "notifications" in data:
                return data["notifications"]
    except (OSError, json.JSONDecodeError):
        pass
    return []


def _save_notifications(notifications: list[dict[str, Any]], user_id: str) -> None:
    """Persist notifications to JSON file."""
    try:
        f = _notifications_file(user_id)
        f.write_text(
            json.dumps(notifications, indent=2, ensure_ascii=False),
            encoding="utf-8",
        )
    except OSError:
        pass


app = FastAPI(
    title="Studaxis API",
    description="Offline-first AI tutoring backend — flashcards, explain, study recommendation",
    version="1.0.0",
)

# Auth routes: /api/auth/signup, /api/auth/login
app.include_router(auth_router)


@app.on_event("startup")
def _startup():
    """Ensure auth DB tables exist on startup. Select hardware-aware model. Ensure Ollama is running and model is pulled."""
    logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(name)s] %(levelname)s: %(message)s")
    init_db()
    model = get_best_model()
    logger.info("Studaxis started; hardware-aware model selected: %s (config: %s)", model, get_config_path_for_log())
    if ensure_ollama_serve():
        if ensure_ollama_model(model):
            logger.info("Ollama ready; model %s available.", model)
        else:
            logger.warning("Ollama running but model %s could not be pulled. Run manually: ollama pull %s", model, model)
    else:
        logger.warning("Ollama not reachable. Install from https://ollama.com and ensure 'ollama serve' is running.")


# CORS: allow local React (Vite 5173), same-origin (8000 default, 6782, 6783)
app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "http://localhost:5173",
        "http://127.0.0.1:5173",
        "http://localhost:3000",
        "http://127.0.0.1:3000",
        "http://localhost:8000",
        "http://127.0.0.1:8000",
        "http://localhost:6782",
        "http://127.0.0.1:6782",
        "http://localhost:6783",
        "http://127.0.0.1:6783",
    ],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Single AI engine instance (lazy init to avoid import-time side effects)
_ai_engine: Optional[AIEngine] = None


def get_ai_engine() -> AIEngine:
    global _ai_engine
    if _ai_engine is None:
        _ai_engine = AIEngine(base_path=str(BASE_PATH))
    return _ai_engine



# ---------------------------------------------------------------------------
# User stats persistence (same schema as preferences.py / Streamlit)
# ---------------------------------------------------------------------------

_DEFAULT_STATS: dict[str, Any] = {
    "user_id": "",
    "last_sync_timestamp": None,
    "streak": {"current": 0, "longest": 0, "last_active_date": None, "milestone_next": 7},
    "quiz_stats": {
        "total_attempted": 0,
        "total_correct": 0,
        "average_score": 0.0,
        "last_quiz_date": None,
        "by_topic": {},
        "total_score_sum": 0,
        "total_max_sum": 0,
        "average_percent": 0,
        "last_score": None,
    },
    "flashcard_stats": {"total_reviewed": 0, "mastered": 0, "due_for_review": 0, "cards": {}},
    "chat_history": [],
    "preferences": {
        "difficulty_level": "Beginner",
        "theme": "light",
        "language": "English",
        "sync_enabled": True,
    },
    "hardware_info": {},
}


def _load_user_stats(user_id: str) -> dict[str, Any]:
    """Load user_stats.json for given user; return defaults on missing/error."""
    try:
        f = _stats_file(user_id)
        if f.exists():
            raw = f.read_text(encoding="utf-8")
            data = json.loads(raw)
            if isinstance(data, dict):
                return data
    except (OSError, json.JSONDecodeError):
        pass
    result = dict(_DEFAULT_STATS)
    prefs = result.get("preferences") or {}
    if not isinstance(prefs, dict):
        prefs = {}
    for key, default in _DEFAULT_STATS["preferences"].items():
        prefs.setdefault(key, default)
    result["preferences"] = prefs
    return result



def _save_user_stats(stats: dict[str, Any], user_id: str) -> None:
    """Persist user stats to per-user directory (atomic write)."""
    try:
        f = _stats_file(user_id)
        tmp = f.with_suffix(".tmp")
        tmp.write_text(json.dumps(stats, indent=2, ensure_ascii=False), encoding="utf-8")
        tmp.replace(f)
    except OSError:
        pass


# ---------------------------------------------------------------------------
# Request/Response models
# ---------------------------------------------------------------------------


class FlashcardGenerateRequest(BaseModel):
    topic_or_chapter: str = Field(..., min_length=1, description="Topic or textbook chapter name")
    input_type: str = Field(default="Topic Name", description="'Topic Name' or 'Textbook Chapter'")
    count: int = Field(default=10, ge=5, le=35, description="Number of flashcards to generate")
    offline_mode: bool = Field(default=True, description="Force local inference")
    user_id: Optional[str] = Field(default=None, description="Optional user identifier")


class FlashcardItem(BaseModel):
    id: str
    topic: str
    front: str
    back: str
    sourceType: Optional[str | list[str]] = None


class FlashcardGenerateResponse(BaseModel):
    cards: list[FlashcardItem]
    topic: str


class FlashcardExplainRequest(BaseModel):
    front: str = Field(..., description="Card question/front")
    back: str = Field(..., description="Card answer/back")
    subject: str = Field(default="General", description="Deck/subject of the card")
    difficulty: str = Field(default="Beginner", description="User profile difficulty level")


class FlashcardExplainResponse(BaseModel):
    text: str
    confidence_score: float = 0.0


class StudyRecommendationRequest(BaseModel):
    topic: str = Field(..., description="Topic or subject for the plan")
    time_budget_minutes: int = Field(default=15, ge=1, le=240, description="Available study time in minutes")
    review_mode: Optional[str] = Field(default="flashcards", description="Context: flashcards, quiz, etc.")
    user_id: Optional[str] = Field(default=None)
    offline_mode: bool = Field(default=True)


class FlashcardRecommendRequest(BaseModel):
    deck_id: str = Field(default="", description="Current deck id (optional for quiz-only mode)")
    subject: str = Field(default="", description="Current subject/topic (optional for quiz-only mode)")
    hard_cards: list[str] = Field(default_factory=list, description="Card fronts marked Hard")
    easy_count: int = Field(default=0)
    hard_count: int = Field(default=0)
    difficulty: str = Field(default="Beginner")
    insights: Optional[dict[str, Any]] = Field(default=None, description="weak_subjects, avg_quiz_score, streak")


class StudyRecommendationResponse(BaseModel):
    text: str
    confidence_score: float = 0.0


class AdaptiveRecommendationResponse(BaseModel):
    """Structured adaptive recommendation (weak topic, action, difficulty)."""
    weak_topic: str = Field(..., description="Weak area to focus on")
    suggested_action: str = Field(..., description="Specific recommendation")
    difficulty_adjustment: str = Field(..., description="Easier / medium / harder")
    text: str = Field(default="", description="Full human-readable summary")
    confidence_score: float = 0.0
    has_data: bool = Field(default=True, description="False when no flashcard or quiz data")


# --- Chat ---
class ChatRequest(BaseModel):
    message: str = Field(..., min_length=1, description="User message")
    is_clarification: bool = Field(default=False, description="True if this is a follow-up clarification")
    task_type: Optional[str] = Field(
        default="chat",
        description="Task type: chat, clarify, explain_topic, quiz_me, flashcards, step_by_step",
    )
    context: Optional[dict[str, Any]] = Field(default=None, description="Optional: difficulty, chat_history, subject, etc.")
    subject: Optional[str] = Field(default=None, description="Selected subject (e.g. General, Maths)")
    textbook_id: Optional[str] = Field(default=None, description="Attached textbook id (filename) for RAG")


class ChatResponse(BaseModel):
    text: str
    confidence_score: float = 0.0
    metadata: Optional[dict[str, Any]] = None


# --- Grade ---
class GradeRequest(BaseModel):
    question_id: str = Field(..., description="Question identifier")
    question: str = Field(..., description="Question text")
    expected_answer: Optional[str] = Field(default=None)
    topic: str = Field(default="General")
    answer: str = Field(..., description="Student answer to grade")
    difficulty: str = Field(default="Beginner")
    rubric: Optional[str] = Field(default=None)
    user_id: Optional[str] = None
    offline_mode: bool = True


class GradeResponse(BaseModel):
    text: str
    confidence_score: float = 0.0
    score: Optional[float] = None
    errors: list[Any] = []
    strengths: list[Any] = []
    remarks: str = ""
    metadata: Optional[dict[str, Any]] = None


# --- Quiz (stub content from Streamlit quiz page) ---
QUIZ_ITEMS: list[dict[str, Any]] = [
    {
        "id": "q1",
        "topic": "Physics",
        "question": "State Newton's second law and explain what each term means.",
        "expected_answer": "Force equals mass times acceleration. F = m * a.",
    },
    {
        "id": "q2",
        "topic": "Biology",
        "question": "What is osmosis in simple terms?",
        "expected_answer": "Movement of water molecules from high concentration to low concentration through a semipermeable membrane.",
    },
    {
        "id": "q3",
        "topic": "Mathematics",
        "question": "Differentiate x^2 and explain the rule used.",
        "expected_answer": "Derivative of x squared is 2x using the power rule.",
    },
]

# Panic Mode exam (Streamlit panic_mode.py)
PANIC_ITEMS: list[dict[str, Any]] = [
    {"id": "p1", "topic": "Physics", "question": "Define momentum and explain one real-world example.", "expected_answer": "Momentum is mass multiplied by velocity."},
    {"id": "p2", "topic": "Biology", "question": "Explain why osmosis is important in plant cells.", "expected_answer": "It helps maintain turgor pressure and water balance."},
    {"id": "p3", "topic": "Mathematics", "question": "What is the derivative of x^3 and which rule applies?", "expected_answer": "3x^2 using the power rule."},
    {"id": "p4", "topic": "Chemistry", "question": "What is pH and what does pH 7 indicate?", "expected_answer": "pH measures acidity/basicity; pH 7 is neutral."},
    {"id": "p5", "topic": "Physics", "question": "State one difference between speed and velocity.", "expected_answer": "Speed is scalar; velocity includes direction."},
]

# Fallback when AI fails or source is "default" — exam always loads (by subject)
FALLBACK_PANIC_QUESTIONS: dict[str, list[dict[str, Any]]] = {
    "Physics": [
        {"id": "fb_p1", "topic": "Physics", "question": "What is Newton's First Law?", "options": ["A. F=ma", "B. Law of inertia", "C. Action-reaction", "D. None"], "correct": 1, "explanation": "An object stays at rest or in motion unless acted upon by an external force."},
        {"id": "fb_p2", "topic": "Physics", "question": "What is the SI unit of force?", "options": ["A. Joule", "B. Watt", "C. Newton", "D. Pascal"], "correct": 2, "explanation": "Force is measured in Newtons (N)."},
    ],
    "Biology": [
        {"id": "fb_b1", "topic": "Biology", "question": "What does DNA stand for?", "options": ["A. Deoxyribonucleic Acid", "B. Deoxyribose Nucleic Acid", "C. Double Nucleic Acid", "D. None of the above"], "correct": 0, "explanation": "DNA = Deoxyribonucleic Acid."},
        {"id": "fb_b2", "topic": "Biology", "question": "Where does photosynthesis occur?", "options": ["A. Mitochondria", "B. Chloroplasts", "C. Nucleus", "D. Ribosomes"], "correct": 1, "explanation": "Chloroplasts contain chlorophyll for photosynthesis."},
    ],
    "General": [
        {"id": "fb_g1", "topic": "General", "question": "What does DNA stand for?", "options": ["A. Deoxyribonucleic Acid", "B. Deoxyribose Nucleic Acid", "C. Double Nucleic Acid", "D. None of the above"], "correct": 0, "explanation": "DNA = Deoxyribonucleic Acid."},
        {"id": "fb_g2", "topic": "General", "question": "What is Newton's First Law?", "options": ["A. F=ma", "B. Law of inertia", "C. Action-reaction", "D. None"], "correct": 1, "explanation": "An object stays at rest or in motion unless acted upon."},
    ],
}

# Flat fallback list for JSON-parse failure (shuffled and returned when Ollama returns invalid JSON)
FALLBACK_PANIC_QUESTIONS_LIST: list[dict[str, Any]] = [
    {"id": "fb1", "type": "mcq", "question": "What is Newton's First Law of Motion?", "options": ["A. F = ma", "B. An object stays at rest or in motion unless acted upon", "C. Every action has an equal and opposite reaction", "D. Energy cannot be created or destroyed"], "correct": 1, "correct_index": 1, "explanation": "Newton's First Law is the law of inertia.", "topic": "Physics"},
    {"id": "fb2", "type": "mcq", "question": "What is the SI unit of force?", "options": ["A. Joule", "B. Watt", "C. Newton", "D. Pascal"], "correct": 2, "correct_index": 2, "explanation": "Force is measured in Newtons (N) where 1N = 1 kg⋅m/s².", "topic": "Physics"},
    {"id": "fb3", "type": "mcq", "question": "Which organelle is the powerhouse of the cell?", "options": ["A. Nucleus", "B. Ribosome", "C. Mitochondria", "D. Golgi body"], "correct": 2, "correct_index": 2, "explanation": "Mitochondria produce ATP through cellular respiration.", "topic": "Biology"},
    {"id": "fb4", "type": "mcq", "question": "What does DNA stand for?", "options": ["A. Deoxyribonucleic Acid", "B. Deoxyribose Nucleic Acid", "C. Double Nucleic Acid", "D. Dinucleotide Acid"], "correct": 0, "correct_index": 0, "explanation": "DNA = Deoxyribonucleic Acid, the molecule of heredity.", "topic": "Biology"},
    {"id": "fb5", "type": "mcq", "question": "What is the chemical formula for water?", "options": ["A. CO2", "B. NaCl", "C. H2O2", "D. H2O"], "correct": 3, "correct_index": 3, "explanation": "Water is H2O — two hydrogen atoms and one oxygen atom.", "topic": "Chemistry"},
]


def _get_fallback_panic_items(subject: str, count: int = 5, question_type: str = "mcq") -> list[dict[str, Any]]:
    """Return fallback panic items for subject when AI fails. Ensures exam always loads."""
    import random
    items = [dict(x) for x in FALLBACK_PANIC_QUESTIONS_LIST]
    random.shuffle(items)
    for i, item in enumerate(items):
        item["id"] = f"fb{i+1}"
        if "correct_index" in item and "correct" not in item:
            item["correct"] = item["correct_index"]
    return items[: min(count, len(items))]


class QuizSubmitRequest(BaseModel):
    answers: list[dict[str, Any]] = Field(..., description="List of {question_id, answer}")
    items: Optional[list[dict[str, Any]]] = Field(default=None, description="Custom quiz items for panic mode grading (when generated from material)")


class QuizSubmitResponse(BaseModel):
    score: float = Field(default=0, description="Total score")
    max_score: float = Field(default=0, description="Max possible score")
    percent: int = Field(default=0, description="Percentage score")
    results: list[dict[str, Any]] = Field(default_factory=list, description="Per-question: question_id, correct, correct_answer, explanation")
    weak_topics_text: Optional[str] = Field(default=None, description="AI weak-topic summary (panic mode only)")
    recommendation_text: Optional[str] = Field(default=None, description="AI study recommendation (panic mode only)")


# ---------------------------------------------------------------------------
# Helpers (mirror Streamlit flashcards.py logic)
# ---------------------------------------------------------------------------


def _extract_json_array(text: str) -> str:
    """Extract a JSON array from model output (strip markdown code blocks)."""
    text = text.strip()
    if text.startswith("```"):
        text = re.sub(r"^```\w*\n?", "", text)
        text = re.sub(r"\n?```\s*$", "", text)
    return text.strip()


def extract_json(text: str) -> str:
    """Extract raw JSON string from Ollama output that may be wrapped in ```json ... ``` or similar."""
    text = (text or "").strip()
    if "```" in text:
        for block in text.split("```"):
            block = block.strip()
            if block.lower().startswith("json"):
                block = block[4:].lstrip()
            if block.startswith("{") or block.startswith("["):
                return block
    if text.startswith("{") or text.startswith("["):
        return text
    for i, ch in enumerate(text):
        if ch in "{[":
            return text[i:]
    return text


def _clean_json(text: str) -> str:
    text = (text or "").strip()

    # strip markdown code blocks
    if "```" in text:
        for block in text.split("```"):
            block = block.strip()
            if block.startswith("json"):
                block = block[4:].strip()
            if block.startswith("{") or block.startswith("["):
                text = block
                break

    # find first JSON character
    start = -1
    for i, ch in enumerate(text):
        if ch in "{[":
            start = i
            break
    if start > 0:
        text = text[start:]

    # find last JSON character
    end = -1
    for i in range(len(text) - 1, -1, -1):
        if text[i] in "}]":
            end = i
            break
    if end >= 0:
        text = text[: end + 1]

    # remove trailing commas (invalid JSON)
    text = re.sub(r",\s*([}\]])", r"\1", text)

    # if it's a bare array, wrap it
    if text.startswith("["):
        text = '{"questions":' + text + "}"

    return text.strip()


def _parse_ai_json(raw: str) -> list:
    """Robustly parse AI response as JSON array. Handles markdown, {\"questions\": [...]}, trailing commas."""
    cleaned = _clean_json(raw)
    # Try parse as object first (e.g. {"questions": [...]})
    try:
        obj = json.loads(cleaned)
        if isinstance(obj, dict) and "questions" in obj:
            q = obj["questions"]
            if isinstance(q, list):
                return q
    except json.JSONDecodeError:
        pass
    start = cleaned.find("[")
    if start == -1:
        raise ValueError("No JSON array found in AI response")
    # Find matching closing bracket to avoid cutting at ] inside a string
    depth = 0
    in_string = None
    escape = False
    i = start
    while i < len(cleaned):
        c = cleaned[i]
        if escape:
            escape = False
            i += 1
            continue
        if c == "\\" and in_string:
            escape = True
            i += 1
            continue
        if in_string:
            if c == in_string:
                in_string = None
            i += 1
            continue
        if c in ('"', "'"):
            in_string = c
            i += 1
            continue
        if c == "[":
            depth += 1
        elif c == "]":
            depth -= 1
            if depth == 0:
                arr_str = cleaned[start : i + 1]
                try:
                    return json.loads(arr_str)
                except json.JSONDecodeError:
                    fixed = re.sub(r",\s*([}\]])", r"\1", arr_str)
                    return json.loads(fixed)
        i += 1
    # Fallback: first [ to last ]
    end = cleaned.rfind("]")
    if end == -1 or end <= start:
        raise ValueError("No JSON array found in AI response")
    arr_str = cleaned[start : end + 1]
    try:
        return json.loads(arr_str)
    except json.JSONDecodeError:
        fixed = re.sub(r",\s*([}\]])", r"\1", arr_str)
        return json.loads(fixed)


def _normalize_cards(raw: list[Any]) -> list[dict[str, Any]]:
    """Convert LLM output to list of cards with id, topic, front, back."""
    out: list[dict[str, Any]] = []
    for i, item in enumerate(raw):
        if not isinstance(item, dict):
            continue
        front = str(item.get("front", item.get("question", ""))).strip() or "?"
        back = str(item.get("back", item.get("answer", ""))).strip() or "—"
        topic = str(item.get("topic", "")).strip() or "General"
        card_id = str(item.get("id", item.get("card_id", str(uuid.uuid4()))))
        out.append({
            "id": card_id,
            "topic": topic,
            "front": front,
            "back": back,
        })
    return out


def _fallback_flashcards_response(subject: str, count: int, source_type: str = "paste") -> FlashcardGenerateResponse:
    """Return minimal valid flashcard response when AI is unavailable (e.g. Ollama down)."""
    n = min(max(2, count), 5)
    placeholders = [
        {"id": "fc1", "topic": subject, "front": "What are the key concepts in this topic?", "back": "Review the main definitions and examples from your notes."},
        {"id": "fc2", "topic": subject, "front": "How would you explain this to someone new to the subject?", "back": "Use simple language and one concrete example."},
        {"id": "fc3", "topic": subject, "front": "What is the most important takeaway?", "back": "Summarize the core idea in one sentence."},
        {"id": "fc4", "topic": subject, "front": "What connections does this have to other topics?", "back": "Consider prerequisites and follow-up concepts."},
        {"id": "fc5", "topic": subject, "front": "What would an exam question on this look like?", "back": "Think of a short question and model answer."},
    ]
    cards = []
    for i in range(n):
        c = dict(placeholders[i])
        c["sourceType"] = source_type
        cards.append(FlashcardItem(**c))
    return FlashcardGenerateResponse(cards=cards, topic=subject or "General")


def _normalize_quiz_items(raw: list[Any], subject: str) -> list[dict[str, Any]]:
    """Convert LLM output to list of quiz items with id, topic, question, expected_answer."""
    out: list[dict[str, Any]] = []
    for i, item in enumerate(raw):
        if not isinstance(item, dict):
            continue
        question = str(item.get("question", "")).strip() or "?"
        expected = str(item.get("expected_answer", item.get("answer", ""))).strip() or ""
        topic = str(item.get("topic", subject)).strip() or subject
        item_id = str(item.get("id", f"pq{i+1}_{uuid.uuid4().hex[:8]}"))
        out.append({
            "id": item_id,
            "topic": topic,
            "question": question,
            "expected_answer": expected,
        })
    return out


# ---------------------------------------------------------------------------
# Endpoints
# ---------------------------------------------------------------------------


@app.get("/api/health")
def health():
    """Liveness/readiness; includes Ollama availability when possible."""
    payload: dict[str, Any] = {"status": "ok", "service": "studaxis-api"}
    try:
        engine = get_ai_engine()
        # Lightweight check: we don't call Ollama here to avoid latency; optional future: ping Ollama
        payload["ollama_available"] = getattr(engine, "config", None) is not None
    except Exception:
        payload["ollama_available"] = False
    return payload


@app.get("/api/ollama/ping")
def ollama_ping():
    """Ping Ollama at localhost:11434. Used by loading screen to wait until local AI is ready.
    Returns friendly message when not ready — no external network, no 500 errors."""
    import requests
    url = os.environ.get("OLLAMA_BASE_URL", "http://localhost:11434")
    base = url.rstrip("/")
    try:
        r = requests.get(f"{base}/", timeout=3)
        ok = r.status_code == 200
        return {"ok": ok, "message": None if ok else "AI is warming up"}
    except requests.exceptions.Timeout:
        return {"ok": False, "message": "AI is warming up. Try again in a moment."}
    except Exception:
        return {"ok": False, "message": "AI is warming up. Ensure Ollama is running."}


def _get_ollama_available_models() -> list[str]:
    """Fetch available model names from Ollama /api/tags. Returns empty list if Ollama unreachable."""
    import requests
    base = os.environ.get("OLLAMA_BASE_URL", "http://localhost:11434").rstrip("/")
    try:
        r = requests.get(f"{base}/api/tags", timeout=5)
        if r.status_code != 200:
            return []
        data = r.json()
        models = data.get("models", [])
        return [m.get("name", "") for m in models if m.get("name")]
    except Exception:
        return []


# Fallback models to try when configured model returns 404 (not found)
# Prefer 7B/8B for 16GB; include 3B for low-RAM
_OLLAMA_FALLBACK_MODELS = [
    "llama3.2:7b", "llama3:8b", "llama3.2:3b-instruct",
    "llama3.2:3b-instruct", "llama3:3b", "llama3.2:3b", "mistral", "llama2",
]


@app.get("/api/ollama/models")
def ollama_models():
    """Return available Ollama models. Used for offline-first notes generation status."""
    models = _get_ollama_available_models()
    configured = get_best_model()
    return {"models": models, "configured": configured, "available": configured in models or any(m.startswith(configured.split(":")[0]) for m in models)}


# ---------------------------------------------------------------------------
# Hardware (Phase 8 — HardwareValidator)
# ---------------------------------------------------------------------------

def _get_hardware_result() -> dict[str, Any]:
    """Run HardwareValidator and return status, message, specs, tips. Safe if psutil missing or validator fails."""
    try:
        from hardware_validator import HardwareValidator
        v = HardwareValidator()
        is_valid, message, specs = v.validate()
        tips = v.get_optimization_tips()
        quantization = v.get_quantization_recommendation()
        status = "block" if not is_valid else ("warn" if tips else "ok")
        return {
            "status": status,
            "message": message,
            "specs": specs,
            "tips": tips,
            "quantization_recommendation": quantization,
            "min_ram_gb": HardwareValidator.MIN_RAM_GB,
            "min_disk_gb": HardwareValidator.MIN_DISK_GB,
            "recommended_ram_gb": HardwareValidator.RECOMMENDED_RAM_GB,
        }
    except Exception as e:
        return {
            "status": "ok",
            "message": "Hardware check unavailable.",
            "specs": {},
            "tips": [],
            "quantization_recommendation": "Q2_K",
            "min_ram_gb": 4.0,
            "min_disk_gb": 2.0,
            "recommended_ram_gb": 6.0,
            "error": str(e),
        }


@app.get("/api/hardware")
def hardware():
    """Hardware check for boot flow: ok / warn / block + specs + tips."""
    return _get_hardware_result()


# ---------------------------------------------------------------------------
# Diagnostics (Settings: Deployment Readiness)
# ---------------------------------------------------------------------------

def _get_sync_readiness(user_id: str) -> tuple[str, str]:
    """Return (sync_state, sync_readiness) from user stats and SyncManager."""
    stats = _load_user_stats(user_id)
    prefs = stats.get("preferences") or {}
    sync_enabled = prefs.get("sync_enabled", True)
    last_sync = stats.get("last_sync_timestamp")

    if not sync_enabled:
        return "Disabled", "Cloud sync is disabled in Settings"

    try:
        from sync_manager import SyncManager
        sm = SyncManager(base_path=str(BASE_PATH), user_id=user_id)
        pending = sm.queue_size
        online = sm.check_connectivity()
        if pending > 0 and online:
            return "Pending", f"{pending} item(s) queued — ready to sync"
        if pending > 0 and not online:
            return "Pending", f"{pending} item(s) queued — waiting for connection"
        if online:
            return "Idle", "Ready to sync"
        return "Offline", "Waiting for connection"
    except Exception:
        pass

    if last_sync:
        return "Idle", "Ready to sync"
    return "Idle", "Ready to sync"


@app.get("/api/diagnostics")
def diagnostics(user_id: str = Depends(get_user_id)):
    """
    Deployment readiness: app version, environment, sync state, last sync.
    Used by Settings Deployment Readiness panel.
    """
    stats = _load_user_stats(user_id)
    prefs = stats.get("preferences") or {}
    sync_enabled = prefs.get("sync_enabled", True)
    last_sync = stats.get("last_sync_timestamp")
    sync_state, sync_readiness = _get_sync_readiness(user_id)

    return {
        "app_version": getattr(app, "version", "1.0.0"),
        "environment": os.environ.get("STUDAXIS_ENV", "Local"),
        "sync_enabled": sync_enabled,
        "sync_state": sync_state,
        "sync_readiness": sync_readiness,
        "last_sync_timestamp": last_sync,
    }


# ---------------------------------------------------------------------------
# Storage file list (Settings: Storage)
# ---------------------------------------------------------------------------

def _format_size(n: int) -> str:
    """Format byte count as human-readable string."""
    if n < 1024:
        return f"{n} B"
    if n < 1024 * 1024:
        return f"{n / 1024:.1f} KB"
    return f"{n / (1024 * 1024):.1f} MB"


def _get_storage_files(user_id: str) -> list[dict[str, Any]]:
    """List storage files for a user (per-user dir) plus shared files."""
    files: list[dict[str, Any]] = []
    user_dir = _user_dir(user_id)

    # Per-user files
    per_user_known = {
        "user_stats.json": "Progress, streaks, quiz stats, chat history, preferences",
        "flashcards.json": "Flashcard decks and SRS data",
        "profile.json": "User profile (name, mode, class code)",
        "users.db": "Auth database (accounts)",
        "sync_queue.json": "Pending sync items",
    }
    for name, desc in per_user_known.items():
        path = user_dir / name
        if path.exists():
            try:
                size = path.stat().st_size
            except OSError:
                size = 0
            extra = ""
            if name == "flashcards.json":
                cards = _load_flashcards(user_id)
                extra = f" — {len(cards)} cards" if cards else ""
            elif name == "user_stats.json":
                stats = _load_user_stats(user_id)
                chat_len = len(stats.get("chat_history") or [])
                if chat_len:
                    extra = f" — {chat_len} chat messages"
            files.append({
                "name": name,
                "size_bytes": size,
                "size_human": _format_size(size),
                "description": desc + extra,
            })

    # Shared files (not per-user)
    shared_known = {
        "profile.json": "User profile (name, mode, class code)",
        "users.db": "Auth database (accounts)",
        "sync_queue.json": "Pending sync items",
    }
    for name, desc in shared_known.items():
        path = DATA_DIR / name
        if path.exists():
            try:
                size = path.stat().st_size
            except OSError:
                size = 0
            files.append({
                "name": name,
                "size_bytes": size,
                "size_human": _format_size(size),
                "description": desc,
            })
        elif name == "sync_queue.json":
            files.append({
                "name": name,
                "size_bytes": 0,
                "size_human": "0 B",
                "description": desc + " — not created yet",
            })

    backups = DATA_DIR / "backups"
    if backups.is_dir():
        try:
            count = sum(1 for _ in backups.iterdir() if _.is_file())
            total = sum(_.stat().st_size for _ in backups.iterdir() if _.is_file())
            files.append({
                "name": "backups/",
                "size_bytes": total,
                "size_human": _format_size(total),
                "description": f"Backup files ({count} file(s))",
            })
        except OSError:
            pass

    return files


@app.get("/api/storage/files")
def storage_files(user_id: str = Depends(get_user_id)):
    """List local storage files for Settings Storage panel."""
    return {"files": _get_storage_files(user_id)}


# ---------------------------------------------------------------------------
# Textbooks (sample_textbooks)
# ---------------------------------------------------------------------------


def _list_textbooks() -> list[dict[str, Any]]:
    """List *.pdf, *.txt, and *.pptx files in sample_textbooks."""
    from datetime import datetime, timezone
    out: list[dict[str, Any]] = []
    if not SAMPLE_TEXTBOOKS_DIR.is_dir():
        return out
    for p in sorted(SAMPLE_TEXTBOOKS_DIR.iterdir()):
        if p.is_file() and p.suffix.lower() in (".pdf", ".txt", ".pptx"):
            subject = p.stem.split("_")[0] if "_" in p.stem else p.stem
            try:
                mtime = p.stat().st_mtime
                uploaded_at = datetime.fromtimestamp(mtime, tz=timezone.utc).isoformat()
            except OSError:
                uploaded_at = ""
            out.append({
                "id": p.name,
                "filename": p.name,
                "name": p.stem,
                "subject": subject,
                "uploaded_at": uploaded_at,
            })
    return out


@app.get("/api/textbooks")
def textbooks_list():
    """List textbooks from data/sample_textbooks (*.pdf, *.txt)."""
    return {"textbooks": _list_textbooks()}


@app.post("/api/textbooks/upload")
def textbooks_upload(file: UploadFile = File(...)):
    """Multipart file upload; save PDF or PPTX to sample_textbooks (shared with Flashcards and AI Chat)."""
    import logging
    log = logging.getLogger("studaxis.textbooks")
    if not file.filename:
        log.warning("[upload] Rejected: no filename")
        raise HTTPException(status_code=400, detail="No filename provided")
    suf = Path(file.filename).suffix.lower()
    if suf not in (".pdf", ".pptx"):
        log.warning("[upload] Rejected: unsupported type %s for %s", suf, file.filename)
        raise HTTPException(status_code=400, detail="Only PDF and PPTX files are accepted")
    SAMPLE_TEXTBOOKS_DIR.mkdir(parents=True, exist_ok=True)
    dest = SAMPLE_TEXTBOOKS_DIR / file.filename
    try:
        content = file.file.read()
        dest.write_bytes(content)
        log.info("[upload] OK: %s -> %s (%d bytes)", file.filename, dest, len(content))
        indexed = True
        try:
            from ai_chat.vector import add_textbook_to_vector_store
            add_textbook_to_vector_store(dest)
        except Exception as idx_err:
            log.warning("[upload] ChromaDB indexing failed for %s: %s", file.filename, idx_err)
            indexed = False
        return {"id": file.filename, "name": Path(file.filename).stem, "indexed": indexed}
    except OSError as e:
        log.error("[upload] FAIL (filesystem): %s: %s", file.filename, e)
        raise HTTPException(status_code=500, detail=f"Could not save file: {e}")
    except Exception as e:
        log.exception("[upload] FAIL: %s: %s", file.filename, e)
        raise HTTPException(status_code=500, detail=f"Upload failed: {e}")


def _extract_text_from_pdf(path: Path) -> str:
    """Extract text from PDF using PyPDF2 or PyPDFLoader."""
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(str(path))
        return "\n".join(p.extract_text() or "" for p in reader.pages)
    except ImportError:
        try:
            from langchain_community.document_loaders import PyPDFLoader
            docs = PyPDFLoader(str(path)).load()
            return "\n".join(d.page_content for d in docs)
        except ImportError:
            raise HTTPException(status_code=501, detail="Install PyPDF2 or langchain-community for PDF extraction")


def _extract_text_from_file(path: Path) -> str:
    """Extract text from txt, pdf, or ppt/pptx."""
    suf = path.suffix.lower()
    if suf == ".txt":
        return path.read_text(encoding="utf-8", errors="replace")
    if suf == ".pdf":
        return _extract_text_from_pdf(path)
    if suf in (".ppt", ".pptx"):
        try:
            from langchain_community.document_loaders import UnstructuredPowerPointLoader
            docs = UnstructuredPowerPointLoader(str(path)).load()
            return "\n".join(d.page_content for d in docs)
        except ImportError:
            raise HTTPException(status_code=501, detail="PPT support requires UnstructuredPowerPointLoader (pip install unstructured)")
    return ""


def _chunks_from_text(text: str, chunk_size: int = 800, overlap: int = 200) -> list[str]:
    """Split text into overlapping chunks for finding relevant content per topic."""
    chunks: list[str] = []
    start = 0
    while start < len(text):
        end = min(start + chunk_size, len(text))
        chunks.append(text[start:end].strip())
        if end >= len(text):
            break
        start += chunk_size - overlap
    return [c for c in chunks if len(c) > 50]


def _find_relevant_chunk_for_topic(content: str, topic: str) -> str:
    """Find the chunk that best matches the topic (contains topic or has high overlap)."""
    chunks = _chunks_from_text(content)
    if not chunks:
        return content[:600] if len(content) > 600 else content
    topic_lower = topic.lower()
    topic_words = set(topic_lower.split())
    best_chunk = chunks[0]
    best_score = 0
    for c in chunks:
        c_lower = c.lower()
        if topic_lower in c_lower or any(w in c_lower for w in topic_words if len(w) > 2):
            overlap = sum(1 for w in topic_words if len(w) > 2 and w in c_lower)
            if overlap > best_score:
                best_score = overlap
                best_chunk = c
    return best_chunk[:800] if len(best_chunk) > 800 else best_chunk


def _generate_single_flashcard_via_ollama(
    topic: str, relevant_chunk: str, subject: str, difficulty: str = "Beginner"
) -> dict[str, Any] | None:
    """Generate one flashcard for a topic using Ollama. Returns {front, back} or None."""
    import requests
    if not relevant_chunk or not str(relevant_chunk).strip():
        context_chunk = (
            f"""Use your knowledge as a {subject} expert to generate questions about: {topic}"""
        )
        no_context_note = (
            f"Since no source material was provided, generate questions based on standard {subject} "
            f"curriculum for {difficulty} level students. "
        )
    else:
        context_chunk = str(relevant_chunk).strip()
        no_context_note = ""

    import random
    variation_seed = random.randint(1, 2_147_483_647)
    variation_instructions = (
        f"VARIATION (seed {variation_seed}): Generate UNIQUE questions. Do NOT repeat previous questions. "
        "Vary question style, phrasing, examples, angles, and difficulty. Be creative and unexpected."
    )
    prompt = f"""You are creating exam-style flashcards for a {difficulty} level {subject} student.

Topic: {topic}
Context from source material: {context_chunk}

Generate 1 flashcard about this topic.
{no_context_note}
{variation_instructions}

RULES:
- Front must be a QUESTION, never just the topic name
- Questions must be specific and testable
- Back must be a direct answer, max 2 sentences
- Each card must test a DIFFERENT aspect of the topic
- Do NOT write "What is {{topic}}?" as the only question

Good question types to use:
  "What causes...?"
  "How does X differ from Y?"
  "What are the main components of...?"
  "Why does...?"
  "Give an example of..."
  "What happens when...?"

Source content to base questions on:
{context_chunk}

Return ONLY a valid JSON array.
No markdown. No backticks. Start with [ end with ]

Each item:
{{
  "front": "specific question about the topic",
  "back": "concise direct answer"
}}"""
    _ollama_base = os.environ.get("OLLAMA_BASE_URL", "http://localhost:11434").rstrip("/")
    _ollama_url = f"{_ollama_base}/api/generate"
    try:
        resp = requests.post(
            _ollama_url,
            json={
                "model": "llama3.2",
                "prompt": prompt,
                "stream": False,
                "options": {"temperature": 0.9, "seed": variation_seed},
            },
            timeout=60,
        )
        resp.raise_for_status()
        raw = (resp.json().get("response") or "").strip()
        cleaned = re.sub(r"```json|```", "", raw).strip()
        obj = {}
        arr_match = re.search(r"\[.*\]", cleaned, re.DOTALL)
        if arr_match:
            try:
                parsed = json.loads(arr_match.group())
                if isinstance(parsed, list) and parsed and isinstance(parsed[0], dict):
                    obj = parsed[0]
                elif isinstance(parsed, dict):
                    obj = parsed
            except json.JSONDecodeError:
                pass
        if not obj:
            match = re.search(r"\{[^{}]*\"front\"[^{}]*\"back\"[^{}]*\}", cleaned, re.DOTALL)
            if not match:
                match = re.search(r"\{.*?\}", cleaned, re.DOTALL)
            if match:
                try:
                    obj = json.loads(match.group())
                except json.JSONDecodeError:
                    pass
        if obj:
            front = str(obj.get("front", "")).strip() or "?"
            back = str(obj.get("back", "")).strip() or "—"
            return {"front": front, "back": back, "topic": topic}
    except Exception as e:
        print(f"[flashcard] Ollama failed for topic {topic}: {e}")
    return None


def _generate_cards_from_textbook(
    content: str, count: int, textbook_id: str
) -> FlashcardGenerateResponse:
    """Topic-based flashcard generation from textbook: extract topics, 1-2 cards per topic."""
    from rag.topic_extractor import extract_dominant_topics

    truncated = content[:12000] if len(content) > 12000 else content
    subject = Path(textbook_id).stem.split("_")[0] if "_" in Path(textbook_id).stem else Path(textbook_id).stem
    if not subject:
        subject = "General"

    topics = extract_dominant_topics(truncated, num_topics=max(5, count // 2))
    if not topics:
        return _generate_cards_from_content(content, count, "textbook", subject, "Beginner")

    cards: list[dict[str, Any]] = []
    for topic in topics:
        if len(cards) >= count:
            break
        relevant = _find_relevant_chunk_for_topic(truncated, topic)
        for _ in range(2):
            if len(cards) >= count:
                break
            card = _generate_single_flashcard_via_ollama(topic, relevant, subject, "Beginner")
            if card:
                card["id"] = str(uuid.uuid4())[:8]
                card["topic"] = topic
                cards.append(card)

    if not cards:
        return _generate_cards_from_content(content, count, "textbook", subject, "Beginner")

    normalized = _normalize_cards([{"id": c.get("id"), "topic": c.get("topic"), "front": c.get("front"), "back": c.get("back")} for c in cards])
    for c in normalized:
        c["sourceType"] = "textbook"
    return FlashcardGenerateResponse(
        cards=[FlashcardItem(**c) for c in normalized],
        topic=subject or "Textbook",
    )


def _generate_cards_topic_aware(
    content: str, count: int, subject: str, source_type: str, difficulty: str = "Beginner"
) -> FlashcardGenerateResponse:
    """Topic extraction + per-topic flashcard generation for URL/file/paste sources."""
    from rag.topic_extractor import extract_dominant_topics

    truncated = content[:8000] if len(content) > 8000 else content
    subj = (subject or "General").strip()
    topics = extract_dominant_topics(truncated, num_topics=max(5, count // 2))
    if not topics:
        return _generate_cards_from_content(content, count, source_type, subj, difficulty)

    cards: list[dict[str, Any]] = []
    for topic in topics:
        if len(cards) >= count:
            break
        relevant = _find_relevant_chunk_for_topic(truncated, topic)
        for _ in range(2):
            if len(cards) >= count:
                break
            card = _generate_single_flashcard_via_ollama(topic, relevant, subj, difficulty)
            if card:
                card["id"] = str(uuid.uuid4())[:8]
                card["topic"] = topic
                cards.append(card)

    if not cards:
        return _generate_cards_from_content(content, count, source_type, subj, difficulty)

    normalized = _normalize_cards([{"id": c.get("id"), "topic": c.get("topic"), "front": c.get("front"), "back": c.get("back")} for c in cards])
    for c in normalized:
        c["sourceType"] = source_type
    return FlashcardGenerateResponse(
        cards=[FlashcardItem(**c) for c in normalized],
        topic=subj or "General",
    )


def _generate_cards_from_content(
    content: str, count: int, source_type: str, subject: str = "General", difficulty: str = "Beginner"
) -> FlashcardGenerateResponse:
    """Generate flashcards from extracted text via AI. Returns fallback cards if AI unavailable."""
    if not content or not content.strip():
        raise HTTPException(status_code=422, detail="No extractable text from source")
    engine = get_ai_engine()
    truncated = content[:12000] if len(content) > 12000 else content
    try:
        response = engine.request(
            task_type=AITaskType.FLASHCARD_GENERATION,
            user_input="",
            context_data={
                "input_type": "Textbook Chapter",
                "topic_or_chapter": "Content-based",
                "count": count,
                "source_content": truncated,
                "subject": subject,
                "difficulty": difficulty,
            },
            offline_mode=True,
            privacy_sensitive=True,
            user_id=None,
        )
    except (ConnectionError, TimeoutError) as e:
        raise HTTPException(
            status_code=503,
            detail=str(e) or "AI is unavailable. Ensure Ollama is running (ollama serve).",
        ) from e
    if response.state in (AIState.FALLBACK_RESPONSE, AIState.ERROR):
        raise HTTPException(
            status_code=503,
            detail=getattr(response, "text", None) or "AI returned no valid response. Ensure Ollama is running.",
        )
    raw_text = _extract_json_array(response.text)
    try:
        parsed = json.loads(raw_text)
    except json.JSONDecodeError as e:
        raise HTTPException(
            status_code=503,
            detail=f"AI returned invalid JSON: {e}. Try again or use a different model.",
        )
    if not isinstance(parsed, list) or not parsed:
        raise HTTPException(status_code=503, detail="AI returned no valid flashcards. Try again.")
    cards = _normalize_cards(parsed)
    if not cards:
        raise HTTPException(status_code=503, detail="AI returned no valid flashcards. Try again.")
    for c in cards:
        c["sourceType"] = source_type
    return FlashcardGenerateResponse(cards=[FlashcardItem(**c) for c in cards], topic="Content-based")


def _generate_mcq_questions_via_ai(content_or_topic: str, subject: str, difficulty: str, count: int) -> list[dict[str, Any]]:
    """Generate MCQ questions via AI from content or topic. Pass content in context_data so the LLM sees it.
    Raises HTTPException 503 if Ollama is unavailable."""
    engine = get_ai_engine()
    context_data: dict[str, Any] = {
        "subject": subject,
        "count": count,
        "difficulty": difficulty,
        "question_format": "mcq",
        "source_content": (content_or_topic or "").strip()[:12000],
    }
    user_input = "Generate multiple choice questions from the provided content below."
    try:
        response = engine.request(
            task_type=AITaskType.QUIZ_GENERATION,
            user_input=user_input,
            context_data=context_data,
            offline_mode=True,
            privacy_sensitive=True,
            user_id=None,
        )
    except (ConnectionError, TimeoutError) as e:
        raise HTTPException(
            status_code=503,
            detail=str(e) or "AI is unavailable. Ensure Ollama is running (ollama serve).",
        ) from e
    raw = (response.text or "").strip()
    if not raw or getattr(response, "state", None) in (AIState.FALLBACK_RESPONSE, AIState.ERROR):
        raise HTTPException(
            status_code=503,
            detail=getattr(response, "text", None) or "AI returned no valid response. Ensure Ollama is running.",
        )
    parsed = None
    for attempt in range(2):
        try:
            parsed = _parse_ai_json(raw)
            break
        except (ValueError, json.JSONDecodeError):
            raw = _clean_json(raw)
            if attempt == 1:
                raise HTTPException(
                    status_code=503,
                    detail="AI returned invalid JSON. Try again or use a different model.",
                )
    if parsed is None:
        parsed = []
    out: list[dict[str, Any]] = []
    for i, item in enumerate(parsed) if isinstance(parsed, list) else []:
        if not isinstance(item, dict):
            continue
        qid = str(item.get("id", f"q{i+1}_{uuid.uuid4().hex[:6]}"))
        text = str(item.get("text", item.get("question", ""))).strip() or "?"
        opts = item.get("options", [])
        if not isinstance(opts, list):
            opts = []
        options = [str(o) for o in opts[:4]]
        while len(options) < 4:
            options.append("(No option)")
        raw_correct = item.get("correct", 0)
        if isinstance(raw_correct, str) and len(raw_correct) == 1:
            letter = raw_correct.upper()
            correct = ord(letter) - ord("A")
            if correct < 0 or correct >= 4:
                correct = 0
        else:
            correct = int(raw_correct) if isinstance(raw_correct, (int, float)) else 0
        if correct < 0 or correct >= 4:
            correct = 0
        explanation = str(item.get("explanation", "")).strip() or ""
        out.append({
            "id": qid,
            "text": text,
            "options": options,
            "correct": correct,
            "explanation": explanation,
        })
    if not out:
        raise HTTPException(status_code=503, detail="AI returned no valid questions. Try again.")
    return out[:count]


def _generate_open_ended_questions_via_ai(content_or_topic: str, subject: str, difficulty: str, count: int) -> list[dict[str, Any]]:
    """Generate open-ended questions via AI from content or topic. Pass content in context_data so the LLM sees it.
    Raises HTTPException 503 if Ollama is unavailable."""
    engine = get_ai_engine()
    context_data: dict[str, Any] = {
        "subject": subject,
        "count": count,
        "difficulty": difficulty,
        "question_format": "open_ended",
        "source_content": (content_or_topic or "").strip()[:12000],
    }
    user_input = "Generate open-ended exam questions from the provided content below."
    try:
        response = engine.request(
            task_type=AITaskType.QUIZ_GENERATION,
            user_input=user_input,
            context_data=context_data,
            offline_mode=True,
            privacy_sensitive=True,
            user_id=None,
        )
    except (ConnectionError, TimeoutError) as e:
        raise HTTPException(
            status_code=503,
            detail=str(e) or "AI is unavailable. Ensure Ollama is running (ollama serve).",
        ) from e
    raw = (response.text or "").strip()
    if not raw or getattr(response, "state", None) in (AIState.FALLBACK_RESPONSE, AIState.ERROR):
        raise HTTPException(
            status_code=503,
            detail=getattr(response, "text", None) or "AI returned no valid response. Ensure Ollama is running.",
        )
    try:
        parsed = _parse_ai_json(raw)
    except (ValueError, json.JSONDecodeError):
        raw = _clean_json(raw)
        try:
            parsed = _parse_ai_json(raw)
        except (ValueError, json.JSONDecodeError):
            raise HTTPException(
                status_code=503,
                detail="AI returned invalid JSON. Try again or use a different model.",
            )
    out: list[dict[str, Any]] = []
    for i, item in enumerate(parsed) if isinstance(parsed, list) else []:
        if not isinstance(item, dict):
            continue
        qid = str(item.get("id", f"q{i+1}_{uuid.uuid4().hex[:6]}"))
        text = str(item.get("text", item.get("question", ""))).strip() or "?"
        sample = str(item.get("sample_answer", item.get("model_answer", item.get("expected_answer", "")))).strip() or ""
        rubric = str(item.get("rubric", "")).strip() or "Assess comprehension and accuracy."
        out.append({
            "id": qid,
            "text": text,
            "sample_answer": sample,
            "rubric": rubric,
        })
    if not out:
        raise HTTPException(status_code=503, detail="AI returned no valid questions. Try again.")
    return out[:count]


def _generate_quiz_from_content(
    content: str,
    subject: str,
    count: int = 5,
    question_type: str = "open_ended",
) -> list[dict[str, Any]]:
    """Generate panic-mode quiz items from extracted text via AI.
    question_type: mcq | open_ended. MCQ returns {id, text, options, correct, explanation}."""
    if not content or not content.strip():
        raise HTTPException(status_code=422, detail="No extractable text from source")
    engine = get_ai_engine()
    truncated = content[:12000] if len(content) > 12000 else content
    qt = question_type if question_type in ("mcq", "open_ended") else "open_ended"

    if qt == "mcq":
        return _generate_mcq_from_content(truncated, subject, count)
    # open_ended
    try:
        response = engine.request(
            task_type=AITaskType.QUIZ_GENERATION,
            user_input="Generate open-ended exam questions from the provided content below.",
            context_data={
                "subject": subject,
                "count": count,
                "source_content": truncated,
            },
            offline_mode=True,
            privacy_sensitive=True,
            user_id=None,
        )
    except (ConnectionError, TimeoutError) as e:
        raise HTTPException(
            status_code=503,
            detail=str(e) or "AI is unavailable. Ensure Ollama is running (ollama serve).",
        ) from e
    if response.state in (AIState.FALLBACK_RESPONSE, AIState.ERROR):
        raise HTTPException(
            status_code=503,
            detail=getattr(response, "text", None) or "AI returned no valid response. Ensure Ollama is running.",
        )
    try:
        parsed = _parse_ai_json(response.text)
    except (ValueError, json.JSONDecodeError) as e:
        raise HTTPException(
            status_code=503,
            detail=f"AI returned invalid JSON: {e}. Try again or use a different model.",
        )
    if not isinstance(parsed, list):
        raise HTTPException(status_code=503, detail="AI returned invalid format. Try again.")
    return _normalize_quiz_items(parsed, subject)


def _generate_mcq_from_content(content: str, subject: str, count: int) -> list[dict[str, Any]]:
    """Generate MCQ questions from extracted content via AI. Raises HTTPException 503 if AI is unavailable."""
    engine = get_ai_engine()
    try:
        response = engine.request(
            task_type=AITaskType.QUIZ_GENERATION,
            user_input="Generate multiple choice questions from the provided content below.",
            context_data={
                "subject": subject,
                "count": count,
                "question_format": "mcq",
                "source_content": content[:12000],
            },
            offline_mode=True,
            privacy_sensitive=True,
            user_id=None,
        )
    except (ConnectionError, TimeoutError) as e:
        raise HTTPException(
            status_code=503,
            detail=str(e) or "AI is unavailable. Ensure Ollama is running (ollama serve).",
        ) from e
    if not (response.text or "").strip() or getattr(response, "state", None) in (AIState.FALLBACK_RESPONSE, AIState.ERROR):
        raise HTTPException(
            status_code=503,
            detail=getattr(response, "text", None) or "AI returned no valid response. Ensure Ollama is running.",
        )
    try:
        parsed = _parse_ai_json(response.text)
    except (ValueError, json.JSONDecodeError) as e:
        raise HTTPException(
            status_code=503,
            detail=f"AI returned invalid JSON: {e}. Try again or use a different model.",
        )
    if not isinstance(parsed, list):
        raise HTTPException(status_code=503, detail="AI returned invalid format. Try again.")
    out: list[dict[str, Any]] = []
    for i, item in enumerate(parsed[:count]):
        if not isinstance(item, dict):
            continue
        qid = str(item.get("id", f"pq{i+1}_{uuid.uuid4().hex[:8]}"))
        text_val = str(item.get("text", item.get("question", ""))).strip() or "?"
        opts = item.get("options", [])
        if not isinstance(opts, list):
            opts = []
        options = [str(o) for o in opts[:4]]
        while len(options) < 4:
            options.append("(No option)")
        raw_correct = item.get("correct", 0)
        if isinstance(raw_correct, str) and len(raw_correct) == 1:
            letter = raw_correct.upper()
            correct = ord(letter) - ord("A")
            if correct < 0 or correct >= 4:
                correct = 0
        else:
            correct = int(raw_correct) if isinstance(raw_correct, (int, float)) else 0
        if correct < 0 or correct >= 4:
            correct = 0
        explanation = str(item.get("explanation", "")).strip() or ""
        out.append({
            "id": qid,
            "question": text_val,
            "text": text_val,
            "topic": subject,
            "options": options,
            "correct": correct,
            "explanation": explanation,
        })
    if not out:
        raise HTTPException(status_code=503, detail="AI returned no valid questions. Try again.")
    return out[:count]


class TextbookGenerateRequest(BaseModel):
    textbook_id: str = Field(..., description="Filename in sample_textbooks")
    chapter: Optional[str] = Field(default=None)
    query: Optional[str] = Field(
        default=None,
        description="Topic or focus, e.g. 'feature selection'. If provided, only relevant chunks from ChromaDB are used as context.",
    )
    count: int = Field(default=10, ge=5, le=35)


@app.post("/api/flashcards/generate/textbook", response_model=FlashcardGenerateResponse)
def flashcards_generate_textbook(req: TextbookGenerateRequest):
    """Generate flashcards from a textbook file. Uses ChromaDB semantic search when possible;
    falls back to file-based extraction if ChromaDB is unavailable or returns no chunks."""
    path = SAMPLE_TEXTBOOKS_DIR / req.textbook_id
    if not path.is_file():
        raise HTTPException(status_code=404, detail=f"Textbook '{req.textbook_id}' not found")

    query = (req.query or req.chapter or "key concepts and definitions").strip()
    retrieved = _get_relevant_chunks_from_chromadb(
        query, k=5, source_filter=req.textbook_id
    )
    if retrieved:
        subject = (
            Path(req.textbook_id).stem.split("_")[0]
            if "_" in Path(req.textbook_id).stem
            else Path(req.textbook_id).stem
        )
        if not subject:
            subject = "General"
        return _generate_cards_from_content(
            retrieved, req.count, "textbook", subject, "Beginner"
        )

    # ChromaDB unavailable or no hits: fall back to file-based flow
    try:
        content = _extract_text_from_file(path)
        if content.strip():
            return _generate_cards_from_textbook(content, req.count, req.textbook_id)
    except Exception:
        pass
    topic = req.chapter or Path(req.textbook_id).stem
    return flashcards_generate(FlashcardGenerateRequest(
        topic_or_chapter=topic,
        input_type="Textbook Chapter",
        count=req.count,
        offline_mode=True,
        user_id=None,
    ))


class WeblinkGenerateRequest(BaseModel):
    url: str = Field(..., min_length=1)
    count: int = Field(default=10, ge=5, le=35)


@app.post("/api/flashcards/generate/weblink", response_model=FlashcardGenerateResponse)
def flashcards_generate_weblink(req: WeblinkGenerateRequest):
    """Fetch URL content, strip HTML, generate flashcards via AI."""
    import requests
    try:
        r = requests.get(req.url.strip(), timeout=15, headers={"User-Agent": "Studaxis/1.0"})
        r.raise_for_status()
        html = r.text
    except Exception as e:
        raise HTTPException(status_code=422, detail=f"Could not fetch URL: {e}")
    text = re.sub(r"<[^>]+>", " ", html)
    text = re.sub(r"\s+", " ", text).strip()
    return _generate_cards_from_content(text, req.count, "weblink")


class GenerateFromUrlRequest(BaseModel):
    url: str = Field(..., min_length=1)
    subject: str = Field(default="General")
    num_cards: int = Field(default=10, ge=5, le=20)
    difficulty: str = Field(default="Beginner")


@app.post("/api/flashcards/generate-from-url", response_model=FlashcardGenerateResponse)
def flashcards_generate_from_url(req: GenerateFromUrlRequest):
    """Scrape URL with BeautifulSoup, topic extraction, smart flashcard generation."""
    url = req.url.strip()
    text = ""
    try:
        import httpx
        from bs4 import BeautifulSoup
    except ImportError:
        httpx = None
        BeautifulSoup = None

    if httpx is not None and BeautifulSoup is not None:
        try:
            response = httpx.get(
                url,
                timeout=10,
                headers={"User-Agent": "Mozilla/5.0 (Windows NT 10.0; rv:91.0) Gecko/20100101 Firefox/91.0"},
            )
            if response.status_code in (401, 403):
                raise HTTPException(
                    status_code=422,
                    detail="This website blocked access. Try pasting the article text directly.",
                )
            response.raise_for_status()
            soup = BeautifulSoup(response.text, "html.parser")
            for tag in soup(["nav", "footer", "script", "style", "header", "aside"]):
                tag.decompose()
            text = soup.get_text(separator="\n", strip=True)[:8000]
        except httpx.HTTPStatusError as e:
            if e.response.status_code in (401, 403):
                raise HTTPException(
                    status_code=422,
                    detail="This website blocked access. Try pasting the article text directly.",
                ) from e
            raise HTTPException(
                status_code=422,
                detail="Could not fetch that URL. Try another link or paste the text directly instead.",
            ) from e
        except (httpx.TimeoutException, httpx.ConnectError, httpx.RequestError):
            raise HTTPException(
                status_code=422,
                detail="Could not fetch that URL. Try another link or paste the text directly instead.",
            )
    else:
        import requests
        try:
            r = requests.get(url, timeout=10, headers={"User-Agent": "Mozilla/5.0"})
            if r.status_code in (401, 403):
                raise HTTPException(
                    status_code=422,
                    detail="This website blocked access. Try pasting the article text directly.",
                )
            r.raise_for_status()
            text = re.sub(r"<[^>]+>", " ", r.text)
            text = re.sub(r"\s+", " ", text).strip()[:8000]
        except requests.exceptions.RequestException:
            raise HTTPException(
                status_code=422,
                detail="Could not fetch that URL. Try another link or paste the text directly instead.",
            )

    if len(text) < 200:
        raise HTTPException(
            status_code=422,
            detail="Not enough content found at that URL. Try a different link.",
        )

    cnt = max(5, min(20, req.num_cards))
    return _generate_cards_topic_aware(text, cnt, req.subject, "weblink", req.difficulty)


class GenerateFromTextRequest(BaseModel):
    text: Optional[str] = Field(default=None)
    paste_text: Optional[str] = Field(default=None, description="Pasted content (alias for text)")
    subject: str = Field(default="General")
    num_cards: int = Field(default=10, ge=5, le=20)
    difficulty: str = Field(default="Beginner")

    def get_text(self) -> str:
        out = (self.paste_text or self.text or "").strip()
        if not out:
            raise ValueError("Either 'text' or 'paste_text' is required")
        return out


@app.post("/api/flashcards/generate-from-text", response_model=FlashcardGenerateResponse)
def flashcards_generate_from_text(req: GenerateFromTextRequest):
    """Paste text: topic extraction + smart flashcard generation."""
    try:
        text = req.get_text()
    except ValueError as e:
        raise HTTPException(status_code=422, detail=str(e))
    if len(text) < 150:
        raise HTTPException(
            status_code=422,
            detail="Please paste at least a paragraph of text to generate flashcards from",
        )
    cnt = max(5, min(20, req.num_cards))
    try:
        return _generate_cards_topic_aware(text[:3000], cnt, req.subject, "paste", req.difficulty)
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(
            status_code=503,
            detail=str(e) or "Flashcard generation failed. Ensure Ollama is running (ollama serve).",
        ) from e


@app.post("/api/flashcards/generate-from-file", response_model=FlashcardGenerateResponse)
def flashcards_generate_from_file(
    file: UploadFile = File(...),
    subject: str = Form("General"),
    num_cards: int = Form(10, ge=5, le=20),
):
    """Multipart: PDF or PPT only. Topic extraction + smart flashcard generation."""
    if not file or not file.filename:
        print("[flashcards] generate-from-file: no file provided")
        raise HTTPException(status_code=400, detail="No file provided")
    suf = Path(file.filename).suffix.lower()
    if suf not in (".pdf", ".ppt", ".pptx"):
        raise HTTPException(status_code=422, detail="Only PDF and PPT files are supported")

    tmp_path = DATA_DIR / "tmp_upload" / Path(file.filename).name
    tmp_path.parent.mkdir(parents=True, exist_ok=True)
    try:
        content = file.file.read()
        tmp_path.write_bytes(content)
        if suf == ".pdf":
            text = _extract_text_from_pdf(tmp_path)
        elif suf in (".ppt", ".pptx"):
            try:
                from pptx import Presentation
                prs = Presentation(str(tmp_path))
                parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            parts.append(shape.text)
                text = "\n".join(parts) if parts else ""
            except ImportError:
                try:
                    from langchain_community.document_loaders import UnstructuredPowerPointLoader
                    docs = UnstructuredPowerPointLoader(str(tmp_path)).load()
                    text = "\n".join(d.page_content for d in docs)
                except ImportError:
                    raise HTTPException(
                        status_code=501,
                        detail="PPT support requires python-pptx or unstructured (pip install python-pptx)",
                    )
        else:
            text = ""
        if not text or len(text.strip()) < 100:
            raise HTTPException(status_code=422, detail="Could not extract enough text from the file")
    finally:
        try:
            tmp_path.unlink(missing_ok=True)
        except OSError:
            pass

    cnt = max(5, min(20, num_cards))
    return _generate_cards_topic_aware(text[:12000], cnt, subject or "General", "file", "Beginner")


@app.post("/api/flashcards/generate/files", response_model=FlashcardGenerateResponse)
def flashcards_generate_files(files: list[UploadFile] = File(...), count: int = Form(10)):
    """Multipart file upload; extract text from txt/pdf/ppt, generate via AI."""
    if not files:
        raise HTTPException(status_code=400, detail="No files provided")
    texts: list[str] = []
    ppt_skipped: list[str] = []
    for uf in files:
        if not uf.filename:
            continue
        suf = Path(uf.filename).suffix.lower()
        if suf not in (".txt", ".pdf", ".ppt", ".pptx"):
            continue
        tmp = Path(uf.filename).name
        try:
            content = uf.file.read()
            tmp_path = DATA_DIR / "tmp_upload" / tmp
            tmp_path.parent.mkdir(parents=True, exist_ok=True)
            tmp_path.write_bytes(content)
            try:
                if suf in (".ppt", ".pptx"):
                    try:
                        from langchain_community.document_loaders import UnstructuredPowerPointLoader
                        docs = UnstructuredPowerPointLoader(str(tmp_path)).load()
                        texts.append("\n".join(d.page_content for d in docs))
                    except ImportError:
                        ppt_skipped.append(uf.filename)
                else:
                    texts.append(_extract_text_from_file(tmp_path))
            finally:
                try:
                    tmp_path.unlink(missing_ok=True)
                except OSError:
                    pass
        except Exception as e:
            raise HTTPException(status_code=422, detail=f"Failed to process {uf.filename}: {e}")
    if ppt_skipped:
        pass  # Already recorded
    combined = "\n\n".join(texts)
    if not combined.strip():
        msg = "No extractable text."
        if ppt_skipped:
            msg += f" PPT support unavailable; skipped: {', '.join(ppt_skipped)}"
        raise HTTPException(status_code=422, detail=msg)
    cnt = max(5, min(35, count))
    return _generate_cards_from_content(combined, cnt, "file")


@app.post("/api/flashcards/generate", response_model=FlashcardGenerateResponse)
def flashcards_generate(req: FlashcardGenerateRequest):
    """
    Generate a deck of flashcards from a topic or chapter name using local AI.
    Returns a list of cards (id, topic, front, back) for the React UI to store or display.
    Returns fallback cards when AI is unavailable.
    """
    topic = req.topic_or_chapter.strip() or "General"
    count = max(5, min(20, req.count))
    engine = get_ai_engine()
    try:
        response = engine.request(
            task_type=AITaskType.FLASHCARD_GENERATION,
            user_input=topic,
            context_data={
                "input_type": req.input_type,
                "topic_or_chapter": topic,
                "count": count,
                "subject": "General",
                "difficulty": "Beginner",
                "source_content": None,
            },
            offline_mode=req.offline_mode,
            privacy_sensitive=True,
            user_id=req.user_id,
        )
    except (ConnectionError, TimeoutError) as e:
        raise HTTPException(
            status_code=503,
            detail=str(e) or "AI is unavailable. Ensure Ollama is running (ollama serve).",
        ) from e

    if response.state in (AIState.FALLBACK_RESPONSE, AIState.ERROR):
        raise HTTPException(
            status_code=503,
            detail=getattr(response, "text", None) or "AI returned no valid response. Ensure Ollama is running.",
        )
    raw_text = _extract_json_array(response.text)
    try:
        parsed = json.loads(raw_text)
    except json.JSONDecodeError as e:
        raise HTTPException(
            status_code=503,
            detail=f"AI returned invalid JSON: {e}. Try again or use a different model.",
        )
    if not isinstance(parsed, list) or not parsed:
        raise HTTPException(status_code=503, detail="AI returned no valid flashcards. Try again.")
    cards = _normalize_cards(parsed)
    if not cards:
        raise HTTPException(status_code=503, detail="AI returned no valid flashcards. Try again.")
    for c in cards:
        c["sourceType"] = "topic"
    return FlashcardGenerateResponse(
        cards=[FlashcardItem(**c) for c in cards],
        topic=topic,
    )


# ---------------------------------------------------------------------------
# Flashcards storage (mirror LocalStorage: flashcards.json)
# ---------------------------------------------------------------------------

def _all_cards_from_decks(decks: list[dict[str, Any]]) -> list[dict[str, Any]]:
    """Flatten all cards from all decks (for backward compat)."""
    cards = []
    for d in decks:
        deck_cards = d.get("cards") or []
        for c in deck_cards:
            c2 = dict(c)
            c2["topic"] = c2.get("topic") or d.get("subject") or d.get("title") or "General"
            cards.append(c2)
    return cards


def _load_flashcards(user_id: str) -> list[dict[str, Any]]:
    """Load all cards by flattening from deck structure."""
    decks = _load_flashcard_decks(user_id)
    return _all_cards_from_decks(decks)


def _save_flashcards(cards: list[dict[str, Any]], user_id: str) -> None:
    """Overwrite by saving as single deck (legacy compat)."""
    deck_id = f"deck_{uuid.uuid4().hex[:12]}"
    normalized = [_normalize_card_for_deck(c) for c in cards]
    deck = {
        "id": deck_id,
        "title": "All Cards",
        "subject": "General",
        "created_at": datetime.now(timezone.utc).strftime("%Y-%m-%d"),
        "cards": normalized,
    }
    _save_flashcard_decks([deck], user_id)


def _recompute_deck_counts(deck: dict[str, Any]) -> None:
    """Update easy_count, hard_count, mastered from cards."""
    cards = deck.get("cards") or []
    easy_count = sum(1 for c in cards if (c.get("ease") or "").lower() == "easy")
    hard_count = sum(1 for c in cards if (c.get("ease") or "").lower() == "hard")
    deck["easy_count"] = easy_count
    deck["hard_count"] = hard_count
    deck["mastered"] = easy_count == len(cards) and len(cards) > 0


def _merge_flashcards(new_cards: list[dict[str, Any]], user_id: str) -> int:
    """
    Merge new cards into decks. Cards with matching id update; new ids append to first deck.
    """
    decks = _load_flashcard_decks(user_id)
    if not decks:
        deck_id = f"deck_{uuid.uuid4().hex[:12]}"
        decks = [{
            "id": deck_id,
            "title": "Imported",
            "subject": "General",
            "created_at": datetime.now(timezone.utc).strftime("%Y-%m-%d"),
            "cards": [],
        }]
    first_deck = decks[0]
    first_cards = first_deck.get("cards") or []
    merged_count = 0
    for c in new_cards:
        cid = c.get("id")
        norm = _normalize_card_for_deck(c)
        if not cid:
            first_cards.append(norm)
            merged_count += 1
            continue
        idx = next((i for i, x in enumerate(first_cards) if (x.get("id") or "") == cid), None)
        if idx is not None:
            first_cards[idx] = {**first_cards[idx], **norm}
            merged_count += 1
        else:
            first_cards.append(norm)
            merged_count += 1
    first_deck["cards"] = first_cards
    _save_flashcard_decks(decks, user_id)
    return merged_count


def _get_due_cards(user_id: str) -> list[dict[str, Any]]:
    """Return cards where next_review <= now or missing."""
    from datetime import datetime, timezone
    now = datetime.now(timezone.utc).isoformat()
    all_cards = _load_flashcards(user_id)
    due = []
    for c in all_cards:
        next_review = c.get("next_review") or ""
        if not next_review or next_review <= now:
            due.append(c)
    return due


@app.get("/api/flashcards")
def flashcards_list(user_id: str = Depends(get_user_id)):
    """Return all decks for the authenticated user."""
    return {"decks": _load_flashcard_decks(user_id)}


@app.get("/api/flashcards/cards")
def flashcards_list_cards(user_id: str = Depends(get_user_id)):
    """Return all cards flattened (backward compat for frontend)."""
    return {"cards": _load_flashcards(user_id)}


def _dashboard_flashcards(user_id: str) -> list[dict[str, Any]]:
    """Transform stored flashcards to dashboard format (id, conceptTitle, content, sourceType)."""
    raw = _load_flashcards(user_id)
    out = []
    for c in raw:
        cid = c.get("id") or str(uuid.uuid4())
        topic = c.get("topic") or "General"
        front = c.get("front") or c.get("question") or ""
        back = c.get("back") or c.get("answer") or ""
        # conceptTitle = topic; content = back (answer). Front is the question.
        concept_title = topic
        content = back or front
        source = c.get("sourceType")
        if source is None:
            source = ["textbook"]
        elif isinstance(source, str):
            source = [source]
        out.append({
            "id": cid,
            "conceptTitle": concept_title,
            "content": content,
            "sourceType": source,
        })
    return out


@app.get("/api/dashboard/flashcards")
def dashboard_flashcards(user_id: str = Depends(get_user_id)):
    """Return flashcards in dashboard format for the authenticated user."""
    return {"cards": _dashboard_flashcards(user_id)}


@app.get("/api/flashcards/due")
def flashcards_due(user_id: str = Depends(get_user_id)):
    """Return cards due for review for the authenticated user."""
    return {"cards": _get_due_cards(user_id)}


class FlashcardsAppendRequest(BaseModel):
    cards: list[dict[str, Any]] = Field(..., description="Cards to append (id, topic, front, back, next_review, etc.)")
    deck_id: Optional[str] = Field(default=None, description="If provided, create/update deck with this id")
    deck_title: Optional[str] = Field(default=None)
    deck_subject: Optional[str] = Field(default=None)


@app.post("/api/flashcards")
def flashcards_append(req: FlashcardsAppendRequest, user_id: str = Depends(get_user_id)):
    """Merge generated cards into storage. If deck_id provided, create/update that deck."""
    if not req.cards:
        return {"ok": True, "appended": 0}
    if req.deck_id:
        decks = _load_flashcard_decks(user_id)
        existing = next((d for d in decks if (d.get("id") or "") == req.deck_id), None)
        normalized = [_normalize_card_for_deck(c) for c in req.cards]
        title = req.deck_title or (req.cards[0].get("topic") if req.cards else "General")
        subject = req.deck_subject or title
        if existing:
            by_id = {c.get("id"): c for c in (existing.get("cards") or [])}
            for c in normalized:
                by_id[c.get("id", "")] = c
            existing["cards"] = list(by_id.values())
            _recompute_deck_counts(existing)
        else:
            decks.insert(0, {
                "id": req.deck_id,
                "title": title,
                "subject": subject,
                "created_at": datetime.now(timezone.utc).strftime("%Y-%m-%d"),
                "cards": normalized,
            })
            _recompute_deck_counts(decks[0])
        _save_flashcard_decks(decks, user_id)
        return {"ok": True, "appended": len(req.cards)}
    n = _merge_flashcards(req.cards, user_id)
    return {"ok": True, "appended": n}


class FlashcardsReplaceRequest(BaseModel):
    cards: list[dict[str, Any]] = Field(..., description="Full list to replace storage with")


@app.put("/api/flashcards")
def flashcards_replace(req: FlashcardsReplaceRequest, user_id: str = Depends(get_user_id)):
    """Replace stored flashcards for the authenticated user."""
    _save_flashcards(req.cards, user_id)
    stats = _load_user_stats(user_id)
    ensure_streak_structure(stats)
    ensure_flashcard_structure(stats)
    update_flashcard_stats_from_cards(stats, req.cards)
    _update_streak(stats)
    _save_user_stats(stats, user_id)
    return {"ok": True, "count": len(req.cards)}


class FlashcardReviewRequest(BaseModel):
    card_id: str = Field(..., description="Card identifier")
    ease: str = Field(..., description="'hard' | 'medium' | 'easy'")
    next_review: str = Field(..., description="Next review date (YYYY-MM-DD)")


@app.post("/api/flashcards/review")
def flashcards_review(req: FlashcardReviewRequest, user_id: str = Depends(get_user_id)):
    """Record a single flashcard review; updates streak and flashcard stats."""
    stats = _load_user_stats(user_id)
    ensure_streak_structure(stats)
    ensure_flashcard_structure(stats)
    mastered = req.ease == "easy"
    update_flashcard_entry(stats, req.card_id, req.ease, req.next_review, mastered=mastered)
    decks = _load_flashcard_decks(user_id)
    for d in decks:
        for c in d.get("cards") or []:
            if (c.get("id") or "") == req.card_id:
                c["ease"] = req.ease
                c["next_review"] = req.next_review
                c["review_count"] = int(c.get("review_count", 0)) + 1
                c["mastered"] = mastered
                _save_flashcard_decks(decks, user_id)
                update_flashcard_stats_from_cards(stats, _all_cards_from_decks(decks))
                _update_streak(stats)
                _save_user_stats(stats, user_id)
                _enqueue_sync(BASE_PATH, user_id, "flashcard_review", {
                    "userId": user_id,
                    "cardId": req.card_id,
                    "ease": req.ease,
                    "nextReview": req.next_review,
                })
                return {"ok": True}
    return {"ok": True}


class FlashcardDeckCreateRequest(BaseModel):
    """Create empty deck. Also accepts full deck save when 'cards' is provided."""
    id: Optional[str] = Field(default=None)
    title: str = Field(..., min_length=1)
    subject: str = Field(default="General")
    source: Optional[str] = Field(default=None)
    card_count: Optional[int] = Field(default=None)
    easy_count: Optional[int] = Field(default=None)
    hard_count: Optional[int] = Field(default=None)
    cards: Optional[list[dict[str, Any]]] = Field(default=None)
    created_at: Optional[str] = Field(default=None)
    last_studied: Optional[str] = Field(default=None)


class FlashcardCardAddRequest(BaseModel):
    deck_id: str = Field(...)
    front: str = Field(..., min_length=1)
    back: str = Field(..., min_length=1)


class FlashcardReviewPatchRequest(BaseModel):
    deck_id: Optional[str] = Field(default=None, description="Deck identifier (required for deck progress)")
    card_id: str = Field(...)
    ease: str = Field(..., description="'hard' | 'medium' | 'easy'")
    next_review: Optional[str] = Field(default=None, description="Next review date (computed if omitted)")


class FlashcardsFromQuizRequest(BaseModel):
    wrong_questions: list[dict[str, Any]] = Field(default_factory=list, description="List of {question_id, text, correct_answer, explanation}")


@app.post("/api/flashcards/from-quiz")
def flashcards_from_quiz(req: FlashcardsFromQuizRequest, user_id: str = Depends(get_user_id)):
    """Create a flashcard deck from wrong quiz questions. Front=question, back=correct_answer."""
    if not req.wrong_questions:
        return {"ok": True, "deck_id": None, "card_count": 0}
    decks = _load_flashcard_decks(user_id)
    deck_id = f"quiz_{uuid.uuid4().hex[:12]}"
    cards: list[dict[str, Any]] = []
    for i, wq in enumerate(req.wrong_questions):
        text = str(wq.get("text", wq.get("question", ""))).strip() or "?"
        correct = str(wq.get("correct_answer", "")).strip() or "—"
        expl = str(wq.get("explanation", "")).strip()
        back = correct
        if expl:
            back = f"{correct}\n\n{expl}"
        cards.append({
            "id": str(wq.get("question_id", f"fc_{uuid.uuid4().hex[:8]}")),
            "front": text,
            "back": back,
            "ease": "medium",
            "next_review": datetime.now(timezone.utc).strftime("%Y-%m-%d"),
            "review_count": 0,
            "mastered": False,
        })
    decks.insert(0, {
        "id": deck_id,
        "title": "Quiz Review — Wrong Answers",
        "subject": "General",
        "created_at": datetime.now(timezone.utc).strftime("%Y-%m-%d"),
        "cards": cards,
        "easy_count": 0,
        "hard_count": len(cards),
        "mastered_count": 0,
    })
    _save_flashcard_decks(decks, user_id)
    return {"ok": True, "deck_id": deck_id, "card_count": len(cards)}


@app.post("/api/flashcards/deck")
def flashcard_create_deck(req: FlashcardDeckCreateRequest, user_id: str = Depends(get_user_id)):
    """Create empty deck or save full deck (when cards provided)."""
    decks = _load_flashcard_decks(user_id)
    # Handle full deck save (cards provided)
    if req.cards and len(req.cards) > 0:
        deck_id = req.id or f"deck_{uuid.uuid4().hex[:12]}"
        normalized = [_normalize_card_for_deck(c) for c in req.cards]
        deck = {
            "id": deck_id,
            "title": req.title,
            "subject": req.subject,
            "source": req.source,
            "card_count": req.card_count or len(normalized),
            "easy_count": req.easy_count or 0,
            "hard_count": req.hard_count or 0,
            "created_at": req.created_at or datetime.now(timezone.utc).strftime("%Y-%m-%d"),
            "last_studied": req.last_studied or datetime.now(timezone.utc).strftime("%Y-%m-%d"),
            "cards": normalized,
        }
        _recompute_deck_counts(deck)
        decks = [d for d in decks if (d.get("id") or "") != deck_id]
        decks.insert(0, deck)
        _save_flashcard_decks(decks, user_id)
        _enqueue_sync(BASE_PATH, user_id, "flashcard_create", {
            "userId": user_id,
            "deckId": deck_id,
            "title": req.title,
            "subject": req.subject,
        })
        return {"ok": True}
    # Create empty deck (legacy)
    deck_id = f"deck_{uuid.uuid4().hex[:12]}"
    deck = {
        "id": deck_id,
        "title": req.title,
        "subject": req.subject,
        "created_at": datetime.now(timezone.utc).strftime("%Y-%m-%d"),
        "cards": [],
    }
    decks.append(deck)
    _save_flashcard_decks(decks, user_id)
    _enqueue_sync(BASE_PATH, user_id, "flashcard_create", {
        "userId": user_id,
        "deckId": deck_id,
        "title": req.title,
        "subject": req.subject,
    })
    return deck


@app.post("/api/flashcards/card")
def flashcard_add_card(req: FlashcardCardAddRequest, user_id: str = Depends(get_user_id)):
    """Add a card to a deck."""
    decks = _load_flashcard_decks(user_id)
    card_id = f"card_{uuid.uuid4().hex[:12]}"
    card = {
        "id": card_id,
        "front": req.front,
        "back": req.back,
        "ease": "medium",
        "next_review": datetime.now(timezone.utc).strftime("%Y-%m-%d"),
        "review_count": 0,
        "mastered": False,
    }
    for d in decks:
        if d.get("id") == req.deck_id:
            d.setdefault("cards", []).append(card)
            _save_flashcard_decks(decks, user_id)
            return card
    raise HTTPException(status_code=404, detail="Deck not found")


@app.patch("/api/flashcards/review")
def flashcard_patch_review(req: FlashcardReviewPatchRequest, user_id: str = Depends(get_user_id)):
    """Update card ease and next_review; update deck easy_count, hard_count, mastered when deck_id provided."""
    stats = _load_user_stats(user_id)
    ensure_streak_structure(stats)
    ensure_flashcard_structure(stats)
    mastered = req.ease == "easy"
    next_review = req.next_review or datetime.now(timezone.utc).strftime("%Y-%m-%d")
    update_flashcard_entry(stats, req.card_id, req.ease, next_review, mastered=mastered)
    decks = _load_flashcard_decks(user_id)
    if req.deck_id:
        deck = next((d for d in decks if (d.get("id") or "") == req.deck_id), None)
        if not deck:
            raise HTTPException(status_code=404, detail="Deck not found")
        target_decks = [deck]
    else:
        target_decks = [d for d in decks if any((c.get("id") or "") == req.card_id for c in (d.get("cards") or []))]
    for deck in target_decks:
        for c in deck.get("cards") or []:
            if (c.get("id") or "") == req.card_id:
                c["ease"] = req.ease
                c["next_review"] = next_review
                c["review_count"] = int(c.get("review_count", 0)) + 1
                c["mastered"] = mastered
                _recompute_deck_counts(deck)
                _save_flashcard_decks(decks, user_id)
                update_flashcard_stats_from_cards(stats, _all_cards_from_decks(decks))
                _update_streak(stats)
                _save_user_stats(stats, user_id)
                _enqueue_sync(BASE_PATH, user_id, "flashcard_review", {
                    "userId": user_id,
                    "deckId": req.deck_id or deck.get("id"),
                    "cardId": req.card_id,
                    "ease": req.ease,
                    "nextReview": next_review,
                })
                return {"ok": True}
    raise HTTPException(status_code=404, detail="Card not found")


@app.delete("/api/flashcards/{card_id}")
def flashcard_delete_card(card_id: str, user_id: str = Depends(get_user_id)):
    """Delete a card by id."""
    decks = _load_flashcard_decks(user_id)
    for d in decks:
        cards = d.get("cards") or []
        for i, c in enumerate(cards):
            if (c.get("id") or "") == card_id:
                cards.pop(i)
                d["cards"] = cards
                _save_flashcard_decks(decks, user_id)
                return {"ok": True}
    raise HTTPException(status_code=404, detail="Card not found")


@app.post("/api/flashcards/explain", response_model=FlashcardExplainResponse)
def flashcards_explain(req: FlashcardExplainRequest):
    """
    Get an AI explanation for a flashcard (front/back). Uses local LLM (Ollama).
    """
    engine = get_ai_engine()
    prompt = (
        f"You are a {req.subject} tutor explaining a flashcard concept to a {req.difficulty} student.\n\n"
        f"Concept: {req.front}\n"
        f"Answer: {req.back}\n\n"
        "Give a clear, concise explanation in 3-4 sentences. Use a simple analogy if helpful. "
        "Do not repeat the question. Just explain why the answer is correct and help the student truly understand it."
    )
    try:
        response = engine.request(
            task_type=AITaskType.FLASHCARD_EXPLANATION,
            user_input=prompt,
            context_data={
                "subject": req.subject,
                "front": req.front,
                "back": req.back,
                "difficulty": req.difficulty,
            },
            offline_mode=True,
            privacy_sensitive=True,
            user_id=None,
        )
    except (ConnectionError, TimeoutError) as e:
        raise HTTPException(status_code=503, detail=str(e))

    return FlashcardExplainResponse(
        text=response.text,
        confidence_score=response.confidence_score,
    )


_NO_DATA_MESSAGE = "Complete a quiz or create flashcards to receive personalized recommendations."


@app.post("/api/flashcards/recommend", response_model=AdaptiveRecommendationResponse)
def flashcards_recommend(req: FlashcardRecommendRequest, user_id: str = Depends(get_user_id)):
    """
    Adaptive study recommendation. NEVER generic plans.
    - If flashcard topic exists: use it (+ quiz data if available)
    - Else if quiz history exists: use quiz avg score, weak topics
    - Else: return has_data=False for UI to show empty message
    """
    stats = _load_user_stats(user_id)
    has_fc = _has_flashcard_topic(req.subject, req.hard_cards)
    has_quiz = _has_quiz_data(stats)

    if not has_fc and not has_quiz:
        return AdaptiveRecommendationResponse(
            weak_topic="",
            suggested_action="",
            difficulty_adjustment="",
            text=_NO_DATA_MESSAGE,
            confidence_score=0.0,
            has_data=False,
        )

    quiz_profile = _get_quiz_profile(stats) if has_quiz else None
    difficulty = req.difficulty or (stats.get("preferences") or {}).get("difficulty_level") or "Beginner"

    if has_fc:
        subject = req.subject.strip() or "General"
        prompt = build_flashcard_based_prompt(
            subject=subject,
            difficulty=difficulty,
            hard_cards=req.hard_cards or [],
            easy_count=req.easy_count or 0,
            hard_count=req.hard_count or 0,
            quiz_profile=quiz_profile,
        )
        fallback_weak = subject
    else:
        profile = quiz_profile or {}
        prompt = build_quiz_only_prompt(difficulty=difficulty, quiz_profile=profile)
        weak_list = profile.get("weak_topics") or []
        fallback_weak = weak_list[0][0] if weak_list else "Your weakest topic"

    engine = get_ai_engine()
    try:
        response = engine.request(
            task_type=AITaskType.STUDY_RECOMMENDATION,
            user_input=prompt,
            context_data={
                "deck_id": req.deck_id,
                "subject": req.subject,
                "hard_cards": req.hard_cards,
                "easy_count": req.easy_count,
                "hard_count": req.hard_count,
                "difficulty": difficulty,
            },
            offline_mode=True,
            privacy_sensitive=True,
            user_id=user_id,
        )
        rec = parse_ai_response(response.text, fallback_weak=fallback_weak)
        return AdaptiveRecommendationResponse(
            weak_topic=rec.weak_topic,
            suggested_action=rec.suggested_action,
            difficulty_adjustment=rec.difficulty_adjustment,
            text=rec.text,
            confidence_score=response.confidence_score,
            has_data=True,
        )
    except (ConnectionError, TimeoutError) as e:
        raise HTTPException(status_code=503, detail=str(e))


@app.post("/api/study/recommendation", response_model=StudyRecommendationResponse)
def study_recommendation(req: StudyRecommendationRequest):
    """
    Get a study plan / recommendation for a topic and time budget (local AI).
    """
    engine = get_ai_engine()
    user_input = f"Suggest a review plan for topic {req.topic}."
    try:
        response = engine.request(
            task_type=AITaskType.STUDY_RECOMMENDATION,
            user_input=user_input,
            context_data={
                "topic": req.topic,
                "time_budget_minutes": req.time_budget_minutes,
                "review_mode": req.review_mode or "flashcards",
            },
            offline_mode=req.offline_mode,
            privacy_sensitive=True,
            user_id=req.user_id,
        )
    except (ConnectionError, TimeoutError) as e:
        raise HTTPException(status_code=503, detail=str(e))

    return StudyRecommendationResponse(
        text=response.text,
        confidence_score=response.confidence_score,
    )


# ---------------------------------------------------------------------------
# Chat
# ---------------------------------------------------------------------------


def _resolve_chat_task_type(req: ChatRequest) -> AITaskType:
    """Map ChatRequest fields to AITaskType."""
    if req.is_clarification:
        return AITaskType.CLARIFY
    raw = (req.task_type or "chat").strip().lower()
    mapping = {
        "chat": AITaskType.CHAT,
        "clarify": AITaskType.CLARIFY,
        "explain_topic": AITaskType.EXPLAIN_TOPIC,
        "quiz_me": AITaskType.QUIZ_ME,
        "flashcards": AITaskType.FLASHCARDS,
        "step_by_step": AITaskType.STEP_BY_STEP,
    }
    return mapping.get(raw, AITaskType.CHAT)


@app.post("/api/chat", response_model=ChatResponse)
def chat(req: ChatRequest, user_id: str = Depends(get_user_id)):
    """Turn-based chat with local LLM. Supports clarification, Explain, Quiz, Flashcards, Step-by-Step."""
    stats = _load_user_stats(user_id)
    ensure_streak_structure(stats)
    _update_streak(stats)
    _save_user_stats(stats, user_id)
    engine = get_ai_engine()
    ctx: dict[str, Any] = dict(req.context) if req.context else {}
    ctx["is_clarification"] = req.is_clarification
    if req.subject is not None:
        ctx["subject"] = req.subject
    if req.textbook_id is not None:
        ctx["textbook_id"] = req.textbook_id
    task_type = _resolve_chat_task_type(req)
    try:
        response = engine.request(
            task_type=task_type,
            user_input=req.message,
            context_data=ctx,
            offline_mode=True,
            privacy_sensitive=True,
            user_id=user_id,
        )
    except (ConnectionError, TimeoutError) as e:
        raise HTTPException(status_code=503, detail=str(e))

    return ChatResponse(
        text=response.text,
        confidence_score=response.confidence_score,
        metadata=response.metadata,
    )


# ---------------------------------------------------------------------------
# Chat history (Layer 2 persistence — backend JSON, restores when localStorage cleared)
# ---------------------------------------------------------------------------


class ChatHistorySession(BaseModel):
    """A saved chat session for history persistence."""
    id: str
    title: str
    messages: list[dict[str, Any]] = Field(default_factory=list)
    timestamp: str
    subject: str = "General"


def _load_chat_history(user_id: str) -> list[dict[str, Any]]:
    """Load chat history for user from JSON file."""
    path = _chat_history_file(user_id)
    if not path.exists():
        return []
    try:
        data = json.loads(path.read_text(encoding="utf-8"))
        return data if isinstance(data, list) else []
    except (json.JSONDecodeError, OSError):
        return []


def _save_chat_history(sessions: list[dict[str, Any]], user_id: str) -> None:
    """Save chat history for user to JSON file."""
    path = _chat_history_file(user_id)
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(sessions, ensure_ascii=False, indent=2), encoding="utf-8")


@app.post("/api/chat/history/save")
def chat_history_save(session: ChatHistorySession, user_id: str = Depends(get_user_id)):
    """Append a saved chat session to the user's chat_history.json."""
    sessions = _load_chat_history(user_id)
    sessions.insert(0, session.model_dump())
    _save_chat_history(sessions, user_id)
    return {"ok": True}


@app.get("/api/chat/history")
def chat_history_get(user_id: str = Depends(get_user_id)):
    """Return array of saved chat sessions for restore when localStorage is empty."""
    sessions = _load_chat_history(user_id)
    return sessions


# ---------------------------------------------------------------------------
# Notifications
# ---------------------------------------------------------------------------


class NotificationPushRequest(BaseModel):
    type: Optional[str] = Field(default="info")
    title: str = Field(..., min_length=1)
    message: Optional[str] = None
    tag: Optional[str] = None
    pinned: Optional[bool] = False
    action: Optional[dict[str, str]] = None


@app.get("/api/notifications")
def notifications_get(user_id: str = Depends(get_user_id)):
    """Return notifications from backend/data/notifications/{user_id}.json."""
    notifs = _load_notifications(user_id)
    return {"notifications": notifs}


@app.post("/api/notifications/push")
def notifications_push(
    req: NotificationPushRequest, user_id: str = Depends(get_user_id)
):
    """Add a notification and return the created item."""
    notifs = _load_notifications(user_id)
    nid = str(uuid.uuid4())
    ts = datetime.now(timezone.utc).isoformat()
    notif: dict[str, Any] = {
        "id": nid,
        "type": req.type or "info",
        "title": req.title,
        "message": req.message,
        "tag": req.tag,
        "pinned": req.pinned or False,
        "read": False,
        "timestamp": ts,
        "action": req.action,
    }
    notifs.insert(0, notif)
    _save_notifications(notifs, user_id)
    return notif


@app.patch("/api/notifications/{notif_id}/read")
def notifications_mark_read(
    notif_id: str, user_id: str = Depends(get_user_id)
):
    """Mark one notification as read."""
    notifs = _load_notifications(user_id)
    for n in notifs:
        if n.get("id") == notif_id:
            n["read"] = True
            _save_notifications(notifs, user_id)
            return {"ok": True}
    raise HTTPException(status_code=404, detail="Notification not found")


@app.delete("/api/notifications/{notif_id}")
def notifications_delete(notif_id: str, user_id: str = Depends(get_user_id)):
    """Remove one notification."""
    notifs = _load_notifications(user_id)
    prev_len = len(notifs)
    notifs[:] = [n for n in notifs if n.get("id") != notif_id]
    if len(notifs) == prev_len:
        raise HTTPException(status_code=404, detail="Notification not found")
    _save_notifications(notifs, user_id)
    return {"ok": True}


@app.delete("/api/notifications/all")
@app.delete("/api/notifications/clear")
def notifications_clear_all(user_id: str = Depends(get_user_id)):
    """Clear all non-pinned notifications. Supports /all and /clear."""
    notifs = _load_notifications(user_id)
    notifs[:] = [n for n in notifs if n.get("pinned")]
    _save_notifications(notifs, user_id)
    return {"ok": True}


@app.post("/api/grade", response_model=GradeResponse)
def grade(req: GradeRequest):
    """Grade subjective/objective answers using AI engine with Red Pen–style feedback."""
    engine = get_ai_engine()
    try:
        response = engine.request(
            task_type=AITaskType.GRADING,
            user_input=req.answer,
            context_data={
                "question_id": req.question_id,
                "question": req.question,
                "expected_answer": req.expected_answer,
                "topic": req.topic,
                "difficulty": req.difficulty,
                "rubric": req.rubric or "[GRADING_RUBRIC_PLACEHOLDER]",
            },
            offline_mode=req.offline_mode,
            privacy_sensitive=True,
            user_id=req.user_id,
        )
    except (ConnectionError, TimeoutError) as e:
        raise HTTPException(status_code=503, detail=str(e))

    # Parse "Score: X/10" from AI response for structured output
    score = None
    import re
    match = re.search(r"Score:\s*(\d+(?:\.\d+)?)\s*/\s*10", response.text, re.IGNORECASE)
    if match:
        score = float(match.group(1))

    return GradeResponse(
        text=response.text,
        confidence_score=response.confidence_score,
        score=score,
        errors=[],  # AI engine returns freeform; optional: parse if needed
        strengths=[],
        remarks="",
        metadata=response.metadata,
    )


# ---------------------------------------------------------------------------
# Quiz
# ---------------------------------------------------------------------------


@app.get("/api/quiz/history")
def quiz_history(user_id: str = Depends(get_user_id)):
    """Return all past quiz results for the authenticated user."""
    return {"results": _load_quiz_history(user_id)}


@app.get("/api/quiz/{quiz_id}")
def quiz_get(quiz_id: str, user_id: str = Depends(get_user_id)):
    """Get quiz content by id. Loads from user file first, then static stubs."""
    loaded = _load_quiz_from_file(user_id, quiz_id)
    if loaded:
        return loaded
    if quiz_id == "default" or quiz_id == "quick":
        return {"id": quiz_id, "items": QUIZ_ITEMS, "title": "Quick Quiz"}
    if quiz_id == "panic":
        return {"id": quiz_id, "items": PANIC_ITEMS, "title": "Panic Mode Exam"}
    for item in QUIZ_ITEMS + PANIC_ITEMS:
        if item["id"] == quiz_id:
            return {"id": quiz_id, "items": [item], "title": f"Quiz: {item['topic']}"}
    raise HTTPException(status_code=404, detail=f"Quiz {quiz_id} not found")


class QuizGenerateRequest(BaseModel):
    source: str = Field(default="topic", description="materials | topic | paste | url (paste/url map to topic_text or separate endpoints)")
    subject: str = Field(default="General")
    source_ids: Optional[list[str]] = Field(default=None, description="Textbook ids when source=materials")
    topic: Optional[str] = Field(default=None, description="Topic name when source=topic (used if topic_text empty)")
    topic_text: Optional[str] = Field(default=None, description="Topic or pasted text when source=topic/paste")
    paste_text: Optional[str] = Field(default=None, description="Pasted content (alias for topic_text when source=paste)")
    query: Optional[str] = Field(default=None, description="Alias for topic_text (some clients send 'query')")
    question_type: str = Field(default="mcq", description="mcq | open_ended")
    num_questions: int = Field(default=10, ge=1, le=20)
    difficulty: str = Field(default="medium")
    url: Optional[str] = Field(default=None)
    textbook_id: Optional[str] = Field(default=None)
    chapter: Optional[str] = Field(default=None)


@app.post("/api/quiz/generate")
def quiz_generate(
    req: QuizGenerateRequest,
    user_id: str = Depends(get_user_id),
):
    """Generate quiz from materials or topic. Saves to data/quizzes/{user_id}/{quiz_id}.json.
    Accepts topic_text, paste_text, query, or topic; uses subject as fallback when all empty."""
    topic = (req.paste_text or req.topic_text or req.query or req.topic or "").strip()
    if req.source == "topic" and not topic:
        raise HTTPException(
            status_code=422,
            detail="topic is required for Quick Topic source",
        )
    if req.source == "materials" and req.source_ids:
        textbook_id = req.source_ids[0]
        query = (
            (req.topic_text or req.topic or req.query or req.subject) or "key concepts"
        ).strip()
        retrieved = _get_relevant_chunks_from_chromadb(
            query, k=6, source_filter=textbook_id
        )
        if retrieved:
            topic = retrieved
        else:
            texts = []
            for tid in req.source_ids[:5]:
                path = SAMPLE_TEXTBOOKS_DIR / tid
                if path.is_file():
                    try:
                        texts.append(_extract_text_from_file(path))
                    except Exception:
                        pass
            if texts:
                topic = "\n\n".join(texts)[:15000]
            if not topic.strip():
                topic = req.subject
    if not topic.strip():
        topic = req.subject or "General Knowledge"
    count = max(1, min(20, req.num_questions))
    difficulty = req.difficulty or "medium"
    subject = req.subject or "General"
    if req.question_type == "open_ended":
        items = _generate_open_ended_questions_via_ai(topic, subject, difficulty, count)
        question_type = "open_ended"
    else:
        items = _generate_mcq_questions_via_ai(topic, subject, difficulty, count)
        question_type = "mcq"
    quiz_id = f"gen_{uuid.uuid4().hex[:12]}"
    payload: dict[str, Any] = {
        "id": quiz_id,
        "title": f"Quiz — {subject}",
        "subject": subject,
        "difficulty": difficulty,
        "question_type": question_type,
        "items": items,
    }
    path = _quiz_file(user_id, quiz_id)
    path.write_text(json.dumps(payload, indent=2, ensure_ascii=False), encoding="utf-8")
    return {"id": quiz_id, "title": payload["title"], "items": items, "subject": subject, "difficulty": difficulty, "question_type": question_type}


class GradeAnswerRequest(BaseModel):
    question: str = Field(...)
    user_answer: str = Field(default="")
    sample_answer: str = Field(default="")
    rubric: str = Field(default="")
    difficulty: str = Field(default="Beginner")


@app.post("/api/quiz/grade-answer")
def quiz_grade_answer(req: GradeAnswerRequest):
    """Grade a single open-ended answer. Returns { score: float, feedback: string }."""
    engine = get_ai_engine()
    try:
        response = engine.request(
            task_type=AITaskType.GRADING,
            user_input=req.user_answer,
            context_data={
                "question": req.question,
                "expected_answer": req.sample_answer,
                "rubric": req.rubric or "Assess comprehension and accuracy.",
                "difficulty": req.difficulty,
            },
            offline_mode=True,
            privacy_sensitive=True,
            user_id=None,
        )
        text = response.text or ""
        score = 5.0
        match = re.search(r"Score:\s*(\d+(?:\.\d+)?)\s*/\s*10", text, re.IGNORECASE)
        if match:
            score = float(match.group(1))
        return {"score": min(10.0, max(0.0, score)), "feedback": text}
    except (ConnectionError, TimeoutError) as e:
        fallback = _local_score(req.user_answer, req.sample_answer)
        return {"score": fallback, "feedback": "AI grading unavailable; scored locally."}


def _local_score(answer: str, expected: str) -> float:
    """Simple deterministic score for progress stats (mirror Streamlit quiz.py)."""
    if not (answer or "").strip():
        return 0.0
    at = set((answer or "").lower().split())
    et = set((expected or "").lower().split())
    if not et:
        return 0.0
    overlap = len(at & et) / len(et)
    return round(min(10.0, max(0.0, overlap * 10)), 1)


class PanicGenerateTextbookRequest(BaseModel):
    subject: str = Field(..., min_length=1)
    textbook_id: str = Field(...)
    chapter: Optional[str] = Field(default=None)
    count: int = Field(default=5, ge=3, le=15)
    question_type: str = Field(default="open_ended", description="mcq | open_ended")


class PanicGenerateWeblinkRequest(BaseModel):
    subject: str = Field(..., min_length=1)
    url: str = Field(..., min_length=1)
    count: int = Field(default=5, ge=3, le=15)
    question_type: str = Field(default="open_ended", description="mcq | open_ended")


@app.post("/api/quiz/panic/generate/textbook")
def panic_generate_textbook(req: PanicGenerateTextbookRequest):
    """Generate panic-mode questions from a textbook. Uses ChromaDB when available; else file extraction."""
    path = SAMPLE_TEXTBOOKS_DIR / req.textbook_id
    if not path.is_file():
        raise HTTPException(status_code=404, detail=f"Textbook '{req.textbook_id}' not found")
    query = (req.chapter or req.subject or "key concepts for exam").strip()
    retrieved = _get_relevant_chunks_from_chromadb(
        query, k=5, source_filter=req.textbook_id
    )
    if retrieved:
        content = retrieved
    else:
        content = _extract_text_from_file(path)
        if not content.strip():
            raise HTTPException(status_code=422, detail="Could not extract text from textbook")
    items = _generate_quiz_from_content(
        content, req.subject, req.count, req.question_type
    )
    return {
        "id": "panic",
        "title": f"Panic Mode — {req.subject}",
        "items": items,
        "question_type": req.question_type,
    }


@app.post("/api/quiz/panic/generate/weblink")
def panic_generate_weblink(req: PanicGenerateWeblinkRequest):
    """Generate panic-mode questions from a web URL. Uses Ollama. Raises 503 if AI is unavailable."""
    url = req.url.strip()
    text = ""
    try:
        import httpx
        from bs4 import BeautifulSoup
    except ImportError:
        httpx = None
        BeautifulSoup = None
    try:
        if httpx is not None and BeautifulSoup is not None:
            try:
                response = httpx.get(
                    url,
                    timeout=20,
                    headers={"User-Agent": "Mozilla/5.0 (Windows NT 10.0; rv:91.0) Gecko/20100101 Firefox/91.0"},
                )
                if response.status_code in (401, 403):
                    raise HTTPException(
                        status_code=422,
                        detail="This website blocked access. Try a different link or paste the text directly.",
                    )
                response.raise_for_status()
                soup = BeautifulSoup(response.text, "html.parser")
                for tag in soup(["nav", "footer", "script", "style", "header", "aside"]):
                    tag.decompose()
                text = soup.get_text(separator="\n", strip=True)[:8000]
            except httpx.HTTPStatusError as e:
                raise HTTPException(status_code=422, detail=f"Could not fetch URL: {e}") from e
            except (httpx.TimeoutException, httpx.ConnectError, httpx.RequestError) as e:
                raise HTTPException(status_code=422, detail=f"Could not fetch URL: {e}") from e
        else:
            import requests as req_lib
            try:
                r = req_lib.get(url, timeout=20, headers={"User-Agent": "Mozilla/5.0 (Windows NT 10.0; rv:91.0) Gecko/20100101 Firefox/91.0"})
                if r.status_code in (401, 403):
                    raise HTTPException(
                        status_code=422,
                        detail="This website blocked access. Try a different link or paste the text directly.",
                    )
                r.raise_for_status()
                text = re.sub(r"<[^>]+>", " ", r.text)
                text = re.sub(r"\s+", " ", text).strip()[:8000]
            except HTTPException:
                raise
            except Exception as e:
                raise HTTPException(status_code=422, detail=f"Could not fetch URL: {e}")
        if len(text) < 150:
            raise HTTPException(
                status_code=422,
                detail="Not enough content found at that URL. Try a different link or use Textbook/Files source.",
            )
        items = _generate_quiz_from_content(
            text, req.subject, req.count, req.question_type
        )
    except HTTPException:
        raise
    return {
        "id": "panic",
        "title": f"Panic Mode — {req.subject}",
        "items": items,
        "question_type": req.question_type,
    }


class QuizGenerateFromUrlRequest(BaseModel):
    url: str = Field(..., min_length=1)
    subject: str = Field(default="General")
    topic_text: Optional[str] = Field(default=None, description="Topic or concept for context")
    num_questions: int = Field(default=10, ge=1, le=20)
    question_type: str = Field(default="mcq", description="mcq | open_ended")
    difficulty: str = Field(default="Beginner")


@app.post("/api/quiz/generate-from-url")
def quiz_generate_from_url(req: QuizGenerateFromUrlRequest, user_id: str = Depends(get_user_id)):
    """Scrape URL, extract content, generate quiz. Returns same format as /api/quiz/generate."""
    url = req.url.strip()
    text = ""
    try:
        import httpx
        from bs4 import BeautifulSoup
    except ImportError:
        httpx = None
        BeautifulSoup = None
    try:
        if httpx is not None and BeautifulSoup is not None:
            try:
                response = httpx.get(
                    url, timeout=15, follow_redirects=True,
                    headers={"User-Agent": "Mozilla/5.0 (Windows NT 10.0; rv:91.0) Gecko/20100101 Firefox/91.0"},
                )
                if response.status_code in (401, 403):
                    raise HTTPException(
                        status_code=422,
                        detail="This website blocked access. Try pasting the article text directly.",
                    )
                response.raise_for_status()
                soup = BeautifulSoup(response.text, "html.parser")
                for tag in soup(["nav", "footer", "script", "style", "header", "aside"]):
                    tag.decompose()
                text = soup.get_text(separator=" ", strip=True)[:4000]
            except HTTPException:
                raise
            except (httpx.HTTPStatusError, httpx.TimeoutException, httpx.ConnectError, httpx.RequestError):
                raise HTTPException(status_code=400, detail="Could not fetch URL")
        else:
            import requests as req_lib
            try:
                r = req_lib.get(url, timeout=15, headers={"User-Agent": "Mozilla/5.0"})
                if r.status_code in (401, 403):
                    raise HTTPException(
                        status_code=422,
                        detail="This website blocked access. Try pasting the article text directly.",
                    )
                r.raise_for_status()
                text = re.sub(r"<[^>]+>", " ", r.text)
                text = re.sub(r"\s+", " ", text).strip()[:4000]
            except HTTPException:
                raise
            except Exception:
                raise HTTPException(status_code=400, detail="Could not fetch URL")
    except HTTPException:
        raise
    except Exception:
        raise HTTPException(status_code=400, detail="Could not fetch URL")
    if len(text) < 200:
        raise HTTPException(
            status_code=422,
            detail="Not enough content found at that URL. Try a different link.",
        )
    return _save_and_return_quiz(user_id, text[:12000], req.subject, req.num_questions, req.question_type, req.difficulty)


class QuizGenerateFromTextRequest(BaseModel):
    text: str = Field(..., min_length=1)
    subject: str = Field(default="General")
    topic_text: Optional[str] = Field(default=None, description="Topic or concept for context")
    num_questions: int = Field(default=10, ge=1, le=20)
    question_type: str = Field(default="mcq", description="mcq | open_ended")
    difficulty: str = Field(default="Beginner")


@app.post("/api/quiz/generate-from-text")
def quiz_generate_from_text(req: QuizGenerateFromTextRequest, user_id: str = Depends(get_user_id)):
    """Generate quiz from pasted text. Returns same format as /api/quiz/generate."""
    text = req.text.strip()
    if len(text) < 150:
        raise HTTPException(
            status_code=422,
            detail="Please paste at least a paragraph of text to generate a quiz from.",
        )
    return _save_and_return_quiz(user_id, text[:3000], req.subject, req.num_questions, req.question_type, req.difficulty)


@app.post("/api/quiz/generate-from-file")
def quiz_generate_from_file(
    file: UploadFile = File(...),
    subject: str = Form("General"),
    topic_text: Optional[str] = Form(None),
    num_questions: int = Form(10, ge=1, le=20),
    question_type: str = Form("mcq"),
    difficulty: str = Form("Beginner"),
    user_id: str = Depends(get_user_id),
):
    """Generate quiz from uploaded PDF or PPT. Returns same format as /api/quiz/generate."""
    if not file or not file.filename:
        raise HTTPException(status_code=400, detail="No file provided")
    suf = Path(file.filename).suffix.lower()
    if suf not in (".pdf", ".ppt", ".pptx"):
        raise HTTPException(status_code=422, detail="Only PDF and PPT files are supported")
    tmp_path = DATA_DIR / "tmp_upload" / Path(file.filename).name
    tmp_path.parent.mkdir(parents=True, exist_ok=True)
    try:
        content = file.file.read()
        tmp_path.write_bytes(content)
        if suf == ".pdf":
            text = _extract_text_from_pdf(tmp_path)
        elif suf in (".ppt", ".pptx"):
            try:
                from pptx import Presentation
                prs = Presentation(str(tmp_path))
                parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            parts.append(shape.text)
                text = "\n".join(parts) if parts else ""
            except ImportError:
                try:
                    from langchain_community.document_loaders import UnstructuredPowerPointLoader
                    docs = UnstructuredPowerPointLoader(str(tmp_path)).load()
                    text = "\n".join(d.page_content for d in docs)
                except ImportError:
                    raise HTTPException(
                        status_code=501,
                        detail="PPT support requires python-pptx (pip install python-pptx)",
                    )
        else:
            text = ""
        if not text or len(text.strip()) < 100:
            raise HTTPException(status_code=422, detail="Could not extract enough text from the file")
    finally:
        try:
            tmp_path.unlink(missing_ok=True)
        except OSError:
            pass
    return _save_and_return_quiz(user_id, text[:12000], subject, num_questions, question_type, difficulty)


def _save_and_return_quiz(
    user_id: str,
    content: str,
    subject: str,
    num_questions: int,
    question_type: str,
    difficulty: str,
) -> Any:
    """Extract content, generate questions, save to quizzes dir, return API response."""
    count = max(1, min(20, num_questions))
    subj = subject or "General"
    if question_type == "open_ended":
        items = _generate_open_ended_questions_via_ai(content, subj, difficulty, count)
    else:
        items = _generate_mcq_questions_via_ai(content, subj, difficulty, count)
    quiz_id = f"gen_{uuid.uuid4().hex[:12]}"
    payload: dict[str, Any] = {
        "id": quiz_id,
        "title": f"Quiz — {subj}",
        "subject": subj,
        "difficulty": difficulty,
        "question_type": question_type,
        "items": items,
    }
    path = _quiz_file(user_id, quiz_id)
    path.write_text(json.dumps(payload, indent=2, ensure_ascii=False), encoding="utf-8")
    return {"id": quiz_id, "title": payload["title"], "items": items, "subject": subj, "difficulty": difficulty, "question_type": question_type}


# ---------------------------------------------------------------------------
# Notes Generator (student-facing — generate study notes from text/textbook)
# ---------------------------------------------------------------------------


class NotesGenerateRequest(BaseModel):
    """Request for /api/notes/generate."""
    source: Optional[str] = Field(default=None, description="paste | topic (optional, for client compatibility)")
    text: Optional[str] = Field(default=None, description="Source text (paste or extracted)")
    paste_text: Optional[str] = Field(default=None, description="Pasted content (alias for text)")
    subject: str = Field(default="General")
    topic: Optional[str] = Field(default=None, description="Optional topic for context")
    style: str = Field(default="summary", description="summary | detailed | revision")

    def get_text(self) -> str:
        out = (self.paste_text or self.text or "").strip()
        if len(out) < 100:
            raise ValueError("At least 100 characters of text or paste_text required")
        return out


class NotesGenerateFromTextbookRequest(BaseModel):
    """Request for /api/notes/generate/textbook."""
    textbook_id: str = Field(..., description="Textbook filename in sample_textbooks")
    subject: str = Field(default="General")
    topic: Optional[str] = Field(default=None)
    style: str = Field(default="summary")


def _call_ollama_generate(ollama_url: str, model: str, prompt: str, req_lib) -> str:
    """Call Ollama /api/generate. Returns response text. Raises on error."""
    resp = req_lib.post(
        ollama_url,
        json={"model": model, "prompt": prompt, "stream": False, "options": {"temperature": 0.5, "num_predict": 2048}},
        timeout=90,
    )
    resp.raise_for_status()
    raw = (resp.json().get("response") or "").strip()
    if not raw:
        raise ValueError("AI returned empty response")
    return raw


def _generate_notes_impl(text: str, subject: str, topic_hint: str, style: str) -> dict[str, Any]:
    """Shared notes generation logic. Returns {generated_text, subject, topic}.
    Uses local Ollama only — fully offline. Tries configured model first, then fallbacks on 404."""
    import requests as req_lib
    style_prompt = {
        "summary": "concise bullet-point summary with key terms in **bold**",
        "detailed": "detailed notes with headings (##), bullets, and key terms",
        "revision": "revision-style notes: short definitions, formulas, and exam tips",
    }.get(style, "concise bullet-point summary")
    prompt = f"""You are a study assistant. Convert the following {subject} content into structured study notes.

Format: Use Markdown.
- Use ## for main headings
- Use bullets (- or *) for lists
- Use **key term** for important terms to remember
- Keep it clear and scannable

Style: {style_prompt}

Content to convert:
---
{text}
---

Output only the notes, no preamble."""

    _ollama_base = os.environ.get("OLLAMA_BASE_URL", "http://localhost:11434").rstrip("/")
    _ollama_url = f"{_ollama_base}/api/generate"
    model = get_best_model()

    # Build list of models to try: configured first, then available, then fallbacks
    available = _get_ollama_available_models()
    models_to_try = [model]
    for m in available:
        if m and m not in models_to_try:
            models_to_try.append(m)
    for m in _OLLAMA_FALLBACK_MODELS:
        if m not in models_to_try:
            models_to_try.append(m)

    last_error: Optional[str] = None
    for try_model in models_to_try:
        try:
            raw = _call_ollama_generate(_ollama_url, try_model, prompt, req_lib)
            return {"generated_text": raw, "subject": subject, "topic": topic_hint or None}
        except req_lib.exceptions.HTTPError as e:
            if e.response is not None and e.response.status_code == 404:
                last_error = f"Model '{try_model}' not found. Run: ollama pull {try_model}"
                continue
            last_error = str(e)
            break
        except req_lib.exceptions.ConnectionError:
            last_error = "Ollama not running. Start with: ollama serve"
            break
        except req_lib.exceptions.Timeout:
            last_error = "AI inference timed out. Try again or use a smaller model."
            break
        except ValueError as ve:
            last_error = str(ve)
            break
        except req_lib.exceptions.RequestException as e:
            last_error = str(e)
            break

    raise HTTPException(
        status_code=503,
        detail=last_error or "Could not generate notes. Ensure Ollama is running (ollama serve) and a model is installed (e.g. ollama pull llama3.2:7b or ollama pull llama3.2:3b-instruct).",
    )


@app.post("/api/notes/generate")
def notes_generate(req: NotesGenerateRequest, user_id: str = Depends(get_user_id)):
    """Generate structured study notes from text using local Ollama."""
    import requests as req_lib
    try:
        text = req.get_text()
    except ValueError as e:
        raise HTTPException(status_code=422, detail=str(e))
    text = text[:8000]
    subject = (req.subject or "General").strip()
    topic_hint = (req.topic or "").strip()
    style = req.style or "summary"
    try:
        return _generate_notes_impl(text, subject, topic_hint, style)
    except req_lib.exceptions.RequestException as e:
        raise HTTPException(status_code=503, detail=str(e))


@app.post("/api/notes/generate/textbook")
def notes_generate_from_textbook(req: NotesGenerateFromTextbookRequest, user_id: str = Depends(get_user_id)):
    """Generate notes from a textbook file. Uses ChromaDB when available; else file extraction."""
    import requests as req_lib
    path = SAMPLE_TEXTBOOKS_DIR / req.textbook_id
    if not path.is_file():
        raise HTTPException(status_code=404, detail=f"Textbook '{req.textbook_id}' not found")
    subject = (req.subject or "General").strip()
    topic_hint = (req.topic or "").strip()
    style = req.style or "summary"
    query = (req.topic or req.subject or "key concepts").strip()
    retrieved = _get_relevant_chunks_from_chromadb(
        query, k=5, source_filter=req.textbook_id
    )
    if retrieved:
        content = retrieved
    else:
        try:
            content = _extract_text_from_file(path)
        except Exception as e:
            raise HTTPException(status_code=422, detail=f"Could not extract text: {e}")
        if not content or len(content.strip()) < 100:
            raise HTTPException(status_code=422, detail="Textbook has insufficient text (min 100 chars)")
        content = content.strip()[:8000]
    try:
        return _generate_notes_impl(content, subject, topic_hint, style)
    except req_lib.exceptions.RequestException as e:
        raise HTTPException(status_code=503, detail=str(e))


@app.post("/api/quiz/panic/generate/files")
def panic_generate_files(
    files: list[UploadFile] = File(...),
    subject: str = Form(...),
    count: int = Form(5),
    question_type: str = Form("open_ended"),
):
    """Generate panic-mode questions from uploaded files. One subject only."""
    if not files:
        raise HTTPException(status_code=400, detail="No files provided")
    texts: list[str] = []
    for uf in files:
        if not uf.filename:
            continue
        suf = Path(uf.filename).suffix.lower()
        if suf not in (".txt", ".pdf", ".ppt", ".pptx"):
            continue
        tmp = Path(uf.filename).name
        try:
            content = uf.file.read()
            tmp_path = DATA_DIR / "tmp_upload" / tmp
            tmp_path.parent.mkdir(parents=True, exist_ok=True)
            tmp_path.write_bytes(content)
            try:
                texts.append(_extract_text_from_file(tmp_path))
            finally:
                try:
                    tmp_path.unlink(missing_ok=True)
                except OSError:
                    pass
        except Exception as e:
            raise HTTPException(status_code=422, detail=f"Failed to process {uf.filename}: {e}")
    combined = "\n\n".join(texts)
    if not combined.strip():
        raise HTTPException(status_code=422, detail="No extractable text from files")
    cnt = max(3, min(15, count))
    qt = question_type if question_type in ("mcq", "open_ended") else "open_ended"
    items = _generate_quiz_from_content(combined, subject, cnt, qt)
    return {
        "id": "panic",
        "title": f"Panic Mode — {subject}",
        "items": items,
        "question_type": qt,
    }


def _quizzes_dir(user_id: str) -> Path:
    """Return per-user quiz directory: data/quizzes/{user_id}/"""
    d = DATA_DIR / "quizzes" / user_id
    d.mkdir(parents=True, exist_ok=True)
    return d


def _quiz_file(user_id: str, quiz_id: str) -> Path:
    """Path for a saved quiz: data/quizzes/{user_id}/{quiz_id}.json"""
    return _quizzes_dir(user_id) / f"{quiz_id}.json"


def _load_quiz_from_file(user_id: str, quiz_id: str) -> dict[str, Any] | None:
    """Load quiz from user's quiz directory if exists."""
    path = _quiz_file(user_id, quiz_id)
    if not path.is_file():
        return None
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError):
        return None


def _quiz_result_file(user_id: str, quiz_id: str, attempt_id: str) -> Path:
    """Path for a single quiz result file."""
    return _quizzes_dir(user_id) / f"{quiz_id}_{attempt_id}_result.json"


def _save_quiz_result(user_id: str, quiz_id: str, result: dict[str, Any]) -> str:
    """Save quiz result to JSON; returns attempt_id."""
    attempt_id = f"{int(datetime.now(timezone.utc).timestamp())}_{uuid.uuid4().hex[:8]}"
    path = _quiz_result_file(user_id, quiz_id, attempt_id)
    path.write_text(json.dumps(result, indent=2, ensure_ascii=False), encoding="utf-8")
    return attempt_id


def _load_quiz_history(user_id: str) -> list[dict[str, Any]]:
    """Load all quiz results from user's quiz directory."""
    d = _quizzes_dir(user_id)
    results = []
    for p in sorted(d.iterdir(), reverse=True):
        if p.is_file() and p.suffix == ".json" and "_result" in p.stem:
            try:
                data = json.loads(p.read_text(encoding="utf-8"))
                results.append(data)
            except (OSError, json.JSONDecodeError):
                pass
    return results


def _resolve_quiz_items(quiz_id: str, items_list: list[dict[str, Any]] | None, user_id: str | None = None) -> list[dict[str, Any]]:
    """Resolve quiz items for grading (from req.items, user file, or static)."""
    if items_list:
        return items_list
    if user_id:
        loaded = _load_quiz_from_file(user_id, quiz_id)
        if loaded and loaded.get("items"):
            return loaded["items"]
    if quiz_id in ("quick", "default"):
        return QUIZ_ITEMS
    if quiz_id == "panic":
        return PANIC_ITEMS
    for item in QUIZ_ITEMS + PANIC_ITEMS:
        if item.get("id") == quiz_id:
            return [item]
    return []


def _score_mcq_answer(user_answer: str, item: dict[str, Any]) -> float:
    """Score MCQ: user_answer can be index (0-3), letter (A-D), or option text."""
    opts = item.get("options") or []
    correct_idx = int(item.get("correct", 0))
    correct_text = opts[correct_idx] if correct_idx < len(opts) else ""
    try:
        idx = int(user_answer.strip())
        return 10.0 if 0 <= idx < len(opts) and idx == correct_idx else 0.0
    except (ValueError, TypeError):
        pass
    ua = (user_answer or "").strip().upper()
    if len(ua) == 1 and "A" <= ua <= "D":
        letter_idx = ord(ua) - ord("A")
        return 10.0 if letter_idx < len(opts) and letter_idx == correct_idx else 0.0
    ua_lower = ua.lower()
    ct = (correct_text or "").strip().lower()
    return 10.0 if ua_lower and ct and ua_lower in ct else _local_score(user_answer, correct_text)


@app.post("/api/quiz/{quiz_id}/submit", response_model=QuizSubmitResponse)
def quiz_submit(quiz_id: str, req: QuizSubmitRequest, user_id: str = Depends(get_user_id)):
    """Submit quiz answers; grade via AI/local and update user stats."""
    engine = get_ai_engine()
    stats = _load_user_stats(user_id)
    ensure_streak_structure(stats)
    quiz_stats = stats.setdefault("quiz_stats", {})
    topic_scores: dict[str, list[float]] = {}
    items_list = _resolve_quiz_items(quiz_id, req.items, user_id)
    total_score = 0.0
    max_score = len(req.answers) * 10.0 if req.answers else 0.0

    for r in req.answers:
        qid = r.get("question_id", "")
        answer_text = str(r.get("user_answer") or r.get("answer", "")).strip()
        score = float(r.get("score", 0))
        if score == 0 and items_list:
            for it in items_list:
                if it.get("id") == qid:
                    if it.get("options"):
                        score = _score_mcq_answer(answer_text, it)
                    else:
                        score = _local_score(answer_text, it.get("expected_answer", it.get("sample_answer", "")))
                    break
        r["score"] = score
        r["answer"] = answer_text
        r["topic"] = r.get("topic") or next(
            (it.get("topic", "General") for it in items_list if it.get("id") == qid),
            "General",
        )
        topic = r["topic"]
        total_score += score
        topic_scores.setdefault(topic, []).append(score)
        total_attempted = int(quiz_stats.get("total_attempted", 0)) + 1
        total_correct = int(quiz_stats.get("total_correct", 0)) + (1 if score >= 6.0 else 0)
        prev_avg = float(quiz_stats.get("average_score", 0.0))
        quiz_stats["total_attempted"] = total_attempted
        quiz_stats["total_correct"] = total_correct
        quiz_stats["average_score"] = round(((prev_avg * (total_attempted - 1)) + score) / total_attempted, 2)
        by_topic = quiz_stats.setdefault("by_topic", {})
        te = by_topic.setdefault(topic, {"attempts": 0, "avg_score": 0.0})
        te["attempts"] = int(te.get("attempts", 0)) + 1
        te["avg_score"] = round(((float(te.get("avg_score", 0)) * (te["attempts"] - 1)) + score) / te["attempts"], 2)
    if max_score > 0:
        _update_quiz_stats(stats, total_score, max_score)
    _update_streak(stats)
    _save_user_stats(stats, user_id)

    weak_topics_text: Optional[str] = None
    recommendation_text: Optional[str] = None
    if topic_scores:
        weak_topics_payload = {
            topic: round(sum(vals) / len(vals), 2) for topic, vals in topic_scores.items()
        }
        engine = get_ai_engine()
        try:
            weak_topic_response = engine.request(
                task_type=AITaskType.WEAK_TOPIC_DETECTION,
                user_input="Identify weak topics from this exam result.",
                context_data={
                    "exam_mode": "panic_mode",
                    "topic_scores": weak_topics_payload,
                    "total_questions": len(items_list),
                },
                offline_mode=True,
                privacy_sensitive=True,
                user_id=None,
            )
            weak_topics_text = weak_topic_response.text
            if weak_topics_text and weak_topics_text.strip():
                rec_response = engine.request(
                    task_type=AITaskType.STUDY_RECOMMENDATION,
                    user_input="Create a post-exam improvement plan.",
                    context_data={
                        "exam_mode": "panic_mode",
                        "topic_scores": weak_topics_payload,
                        "weak_topics_summary": weak_topics_text,
                        "study_time_minutes": 20,
                    },
                    offline_mode=True,
                    privacy_sensitive=True,
                    user_id=None,
                )
                recommendation_text = rec_response.text
        except (ConnectionError, TimeoutError):
            pass
        if not weak_topics_text:
            weak_topics_text = "AI unavailable for weak-topic analysis."
        if not recommendation_text:
            recommendation_text = "Complete more quizzes and review weak areas from your stats."

    subject = next((r.get("topic", "General") for r in req.answers), "General") if req.answers else "General"
    percent = int(round((total_score / max_score * 100))) if max_score > 0 else 0
    quiz_meta = _load_quiz_from_file(user_id, quiz_id) or {}
    qtype = quiz_meta.get("question_type", "open_ended")

    def _correct_answer_and_explanation(qid: str, it: dict[str, Any], scr: float) -> tuple[str, str]:
        opts = it.get("options")
        if opts:
            idx = int(it.get("correct", 0))
            correct = opts[idx] if idx < len(opts) else ""
            expl = str(it.get("explanation", "")).strip()
            return correct, expl
        return (
            str(it.get("sample_answer", it.get("expected_answer", ""))).strip(),
            str(it.get("explanation", "")).strip(),
        )

    results_out: list[dict[str, Any]] = []
    for r in req.answers:
        qid = r.get("question_id", "")
        scr = float(r.get("score", 0))
        correct = scr >= 6.0
        correct_ans, explanation = "", ""
        for it in items_list:
            if it.get("id") == qid:
                correct_ans, explanation = _correct_answer_and_explanation(qid, it, scr)
                break
        if not correct and r.get("feedback"):
            explanation = r.get("feedback", explanation)
        results_out.append({
            "question_id": qid,
            "correct": correct,
            "score": scr,
            "correct_answer": correct_ans,
            "explanation": explanation,
        })

    answers_payload = [
        {
            "question_id": r.get("question_id", ""),
            "user_answer": r.get("answer", ""),
            "correct": float(r.get("score", 0)) >= 6.0,
            "score": float(r.get("score", 0)),
            "correct_answer": next((ro["correct_answer"] for ro in results_out if ro["question_id"] == r.get("question_id")), ""),
            "explanation": next((ro["explanation"] for ro in results_out if ro["question_id"] == r.get("question_id")), r.get("feedback", "")),
        }
        for r in req.answers
    ]
    result_payload = {
        "quiz_id": quiz_id,
        "completed_at": datetime.now(timezone.utc).isoformat(),
        "score": total_score,
        "max_score": max_score,
        "percent": percent,
        "subject": subject,
        "question_type": qtype,
        "answers": answers_payload,
    }
    _save_quiz_result(user_id, quiz_id, result_payload)
    _enqueue_sync(BASE_PATH, user_id, "quiz_result", {
        "userId": user_id,
        "quizId": quiz_id,
        "result": result_payload,
    })
    # Enqueue recordQuizAttempt for AppSync (local-first: save first, then background sync)
    _enqueue_quiz_sync_for_submit(user_id, quiz_id, total_score, max_score, len(items_list), subject, qtype)

    if req.answers:
        _enqueue_panic_quiz_for_sync(req.answers, len(items_list), user_id)

    return {
        "score": total_score,
        "max_score": max_score,
        "percent": percent,
        "results": results_out,
        "weak_topics_text": weak_topics_text,
        "recommendation_text": recommendation_text,
    }


@app.get("/api/student/assignments")
def student_assignments(class_code: str = "", user_id: str = Depends(get_user_id)):
    """Return assignments for class. Used when profile.class_code is set (teacher_linked)."""
    items = _load_assignments(class_code)
    return [{"id": a.get("id"), "quiz_id": a.get("quiz_id"), "title": a.get("title"), "due_date": a.get("due_date"), "assigned_at": a.get("assigned_at"), "status": a.get("status", "pending")} for a in items]


class AssignmentCompleteRequest(BaseModel):
    assignment_id: str = Field(...)
    score: float = Field(default=0)
    completed_at: str = Field(default_factory=lambda: datetime.now(timezone.utc).isoformat())


@app.post("/api/student/assignment-complete")
def student_assignment_complete(req: AssignmentCompleteRequest, user_id: str = Depends(get_user_id)):
    """Mark assignment as completed. If offline, enqueue for sync."""
    try:
        p = load_profile_for_user(user_id)
        cc = (p.class_code if p else "") or ""
        if not cc:
            return {"ok": True}
        items = _load_assignments(cc)
        for a in items:
            if a.get("id") == req.assignment_id:
                a["status"] = "completed"
                a["completed_at"] = req.completed_at
                a["score"] = req.score
                _save_assignments(cc, items)
                return {"ok": True}
    except Exception:
        pass
    _enqueue_sync(BASE_PATH, user_id, "assignment_complete", {
        "userId": user_id,
        "assignment_id": req.assignment_id,
        "score": req.score,
        "completed_at": req.completed_at,
    })
    return {"ok": True}


def _teachers_dir() -> Path:
    """Data directory for teacher records (onboarding, class creation)."""
    d = DATA_DIR / "teachers"
    d.mkdir(parents=True, exist_ok=True)
    return d


def _teacher_file(class_code: str) -> Path:
    """Path to teacher record for a class. Keyed by class_code."""
    safe = re.sub(r"[^a-zA-Z0-9_-]", "_", (class_code or "").strip())[:64]
    return _teachers_dir() / f"{safe or 'default'}.json"


def _load_teacher(class_code: str) -> dict[str, Any] | None:
    """Load teacher record for class. Returns None if not found."""
    try:
        path = _teacher_file(class_code)
        if path.is_file():
            return json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError):
        pass
    return None


def _save_teacher(class_code: str, data: dict[str, Any]) -> None:
    """Save teacher record for class."""
    try:
        path = _teacher_file(class_code)
        path.write_text(json.dumps(data, indent=2, ensure_ascii=False), encoding="utf-8")
    except OSError:
        pass


class TeacherOnboardRequest(BaseModel):
    """Teacher onboarding data from StudaxisTeacherDashboard.jsx flow."""
    name: str = Field(..., min_length=1)
    email: str = Field(..., min_length=1)
    subject: str = Field(default="")
    grade: str = Field(default="")
    school: str = Field(default="")
    city: str = Field(default="")
    board: str = Field(default="")
    className: str = Field(..., min_length=1)
    classCode: str = Field(..., min_length=3)
    numStudents: str = Field(default="")


@app.post("/api/teacher/onboard")
def teacher_onboard(req: TeacherOnboardRequest):
    """
    Register teacher + create first class. Public endpoint (no auth).
    Persists to data/teachers/{class_code}.json for backend assignment routing.
    Called by teacher dashboard after onboarding completion.
    """
    cc = req.classCode.strip()
    existing = _load_teacher(cc)
    if existing:
        logger.info("Teacher class code %s already registered; updating", cc)
    payload = {
        "name": req.name,
        "email": req.email.strip().lower(),
        "subject": req.subject,
        "grade": req.grade,
        "school": req.school,
        "city": req.city,
        "board": req.board,
        "className": req.className,
        "classCode": cc,
        "numStudents": req.numStudents,
        "onboarded_at": datetime.now(timezone.utc).isoformat(),
    }
    _save_teacher(cc, payload)
    return {"ok": True, "classCode": cc}


@app.get("/api/teacher/lookup")
def teacher_lookup(classCode: str = ""):
    """
    Look up teacher by class code. Public endpoint (no auth).
    Used by teacher dashboard Login when returning teachers sign in.
    Returns 404 if class code not found.
    """
    cc = (classCode or "").strip()
    if not cc or len(cc) < 3:
        raise HTTPException(status_code=400, detail="classCode is required (min 3 chars)")
    teacher = _load_teacher(cc)
    if not teacher:
        raise HTTPException(status_code=404, detail="Class code not found")
    return teacher


class TeacherAuthRequest(BaseModel):
    """Login payload from teacher dashboard UI. Matches exact fields: classCode (required), teacherId (optional)."""
    classCode: str = Field(..., min_length=3, description="Class code e.g. CS101, PHYS11A")
    teacherId: str | None = Field(default=None, description="Optional, for verification/offline fallback")


def _create_teacher_jwt(class_code: str, teacher_id: str) -> str:
    """Create JWT for authenticated teacher. Uses same secret as auth_routes for consistency."""
    import jwt
    from auth_routes import JWT_ALGORITHM, JWT_EXPIRY_HOURS
    from auth_routes import JWT_SECRET
    from datetime import timedelta
    payload = {
        "sub": class_code,
        "classCode": class_code,
        "teacherId": teacher_id,
        "role": "teacher",
        "exp": datetime.now(timezone.utc) + timedelta(hours=JWT_EXPIRY_HOURS),
        "iat": datetime.now(timezone.utc),
    }
    return jwt.encode(payload, JWT_SECRET, algorithm=JWT_ALGORITHM)


@app.post("/api/teacher/auth")
def teacher_auth(req: TeacherAuthRequest):
    """
    Authenticate teacher by class code. Verifies classCode exists in teacher store.
    Returns JWT + teacher object. Matches Login UI fields: classCode (required), teacherId (optional).
    """
    cc = (req.classCode or "").strip().upper()
    if len(cc) < 3:
        raise HTTPException(status_code=400, detail="classCode is required (min 3 chars)")
    teacher = _load_teacher(cc)
    if not teacher:
        raise HTTPException(status_code=401, detail="Class code not found")
    teacher_id = teacher.get("teacherId") or teacher.get("classCode") or cc
    if req.teacherId and req.teacherId.strip():
        tid = req.teacherId.strip()
        if teacher.get("email", "").lower() == tid.lower() or teacher.get("teacherId") == tid:
            teacher_id = tid
        else:
            teacher_id = tid
    token = _create_teacher_jwt(cc, teacher_id)
    teacher_response = {
        "teacherId": teacher_id,
        "name": teacher.get("name", ""),
        "email": teacher.get("email", ""),
        "subject": teacher.get("subject", ""),
        "grade": teacher.get("grade", ""),
        "school": teacher.get("school", ""),
        "city": teacher.get("city", ""),
        "board": teacher.get("board", ""),
        "className": teacher.get("className", ""),
        "classCode": teacher.get("classCode", cc),
        "numStudents": teacher.get("numStudents", ""),
    }
    return {"access_token": token, "token_type": "bearer", "teacher": teacher_response}


class TeacherAssignQuizRequest(BaseModel):
    class_code: str = Field(..., min_length=1)
    quiz_id: str = Field(..., min_length=1)
    title: str = Field(..., min_length=1)
    due_date: str = Field(default="", description="YYYY-MM-DD or empty")


@app.post("/api/teacher/assign-quiz")
def teacher_assign_quiz(req: TeacherAssignQuizRequest, user_id: str = Depends(get_user_id)):
    """Create assignment for class. Saves to data/assignments/{class_code}.json"""
    items = _load_assignments(req.class_code)
    aid = str(uuid.uuid4())[:8]
    items.append({
        "id": aid,
        "quiz_id": req.quiz_id,
        "title": req.title,
        "due_date": req.due_date or "",
        "assigned_at": datetime.now(timezone.utc).isoformat(),
        "status": "pending",
    })
    _save_assignments(req.class_code, items)
    return {"ok": True, "assignment_id": aid}


def _enqueue_quiz_sync_for_submit(
    user_id: str, quiz_id: str, total_score: float, max_score: float,
    total_questions: int, subject: str, question_type: str,
) -> None:
    """Queue quiz attempt for AWS AppSync sync. Local-first: save to user_stats first, then background sync."""
    try:
        if not user_id or user_id == "anonymous":
            return
        prefs = _load_user_stats(user_id).get("preferences") or {}
        if not prefs.get("sync_enabled", True):
            return
        percent = int(round((total_score / max_score * 100))) if max_score > 0 else 0
        from sync_manager import SyncManager
        sm = SyncManager(base_path=str(BASE_PATH), user_id=user_id)
        sm.enqueue_quiz_sync(
            user_id=user_id,
            quiz_id=quiz_id,
            score=min(100, max(0, percent)),
            total_questions=total_questions,
            subject=subject or "General",
            difficulty="Medium",
        )
    except Exception:
        pass  # Sync is best-effort; do not fail the request


def _enqueue_panic_quiz_for_sync(results: list[dict[str, Any]], total_questions: int, user_id: str) -> None:
    """Queue panic mode quiz attempt for AWS AppSync sync. No-op if sync disabled or anonymous."""
    try:
        if not user_id or user_id == "anonymous":
            return
        prefs = _load_user_stats(user_id).get("preferences") or {}
        if not prefs.get("sync_enabled", True):
            return
        avg = sum(float(r.get("score", 0)) for r in results) / len(results) if results else 0.0
        score_pct = int(round(avg * 10))  # 0–10 scale → 0–100 for AWS
        from sync_manager import SyncManager
        sm = SyncManager(base_path=str(BASE_PATH))
        sm.enqueue_quiz_sync(
            user_id=user_id,
            quiz_id="panic",
            score=min(100, max(0, score_pct)),
            total_questions=total_questions,
            subject="Panic Mode",
            difficulty="Medium",
        )
    except Exception:
        pass  # Sync is best-effort; do not fail the request


class PanicGradeOneRequest(BaseModel):
    question_id: str = Field(...)
    answer: str = Field(default="")
    question: str = Field(...)
    topic: str = Field(default="General")
    expected_answer: str = Field(default="")
    question_type: str = Field(default="open_ended", description="mcq | open_ended")
    options: Optional[list[str]] = Field(default=None, description="For MCQ: option strings")
    correct: Optional[int] = Field(default=None, description="For MCQ: correct option index 0-3")


class PanicFinalizeRequest(BaseModel):
    results: list[dict[str, Any]] = Field(..., description="Per-question results from grade-one")
    items: list[dict[str, Any]] = Field(..., description="Quiz items for stats")


@app.post("/api/quiz/panic/grade-one")
def panic_grade_one(req: PanicGradeOneRequest):
    """Grade a single panic-mode question. For MCQ uses _score_mcq_answer; for open-ended uses AI + _local_score fallback."""
    engine = get_ai_engine()
    if req.question_type == "mcq" and req.options and req.correct is not None:
        item = {
            "options": req.options,
            "correct": req.correct,
            "expected_answer": req.expected_answer or (req.options[req.correct] if req.correct < len(req.options) else ""),
        }
        score = _score_mcq_answer(req.answer, item)
        feedback_text = "Correct!" if score >= 10.0 else (
            f"The correct answer is: {item.get('expected_answer', '')}"
        )
        return {
            "question_id": req.question_id,
            "score": score,
            "feedback": feedback_text,
            "topic": req.topic,
        }
    score = _local_score(req.answer, req.expected_answer)
    feedback_text = "Grading unavailable."
    try:
        grading = engine.request(
            task_type=AITaskType.GRADING,
            user_input=req.answer,
            context_data={
                "question_id": req.question_id,
                "question": req.question,
                "topic": req.topic,
                "expected_answer": req.expected_answer,
                "difficulty": "Beginner",
                "rubric": "[GRADING_RUBRIC_PLACEHOLDER]",
            },
            offline_mode=True,
            privacy_sensitive=True,
        )
        if grading.state in (AIState.TIMEOUT, AIState.ERROR, AIState.FALLBACK_RESPONSE):
            feedback_text = "AI grading timed out; scored locally."
        else:
            feedback_text = grading.text or "No feedback."
    except (ConnectionError, TimeoutError):
        feedback_text = "AI grading timed out; scored locally."
    return {
        "question_id": req.question_id,
        "score": score,
        "feedback": feedback_text,
        "topic": req.topic,
    }


@app.post("/api/quiz/panic/finalize")
def panic_finalize(req: PanicFinalizeRequest, user_id: str = Depends(get_user_id)):
    """Update stats from pre-graded results and return weak topics + recommendation. Falls back on timeout."""
    stats = _load_user_stats(user_id)
    ensure_streak_structure(stats)
    quiz_stats = stats.setdefault("quiz_stats", {})
    topic_scores: dict[str, list[float]] = {}
    total_score = 0.0
    max_score = len(req.results) * 10.0 if req.results else 0.0
    for r in req.results:
        topic = r.get("topic", "General")
        score = float(r.get("score", 0))
        total_score += score
        topic_scores.setdefault(topic, []).append(score)
        total_attempted = int(quiz_stats.get("total_attempted", 0)) + 1
        total_correct = int(quiz_stats.get("total_correct", 0)) + (1 if score >= 6.0 else 0)
        prev_avg = float(quiz_stats.get("average_score", 0.0))
        quiz_stats["total_attempted"] = total_attempted
        quiz_stats["total_correct"] = total_correct
        quiz_stats["average_score"] = round(((prev_avg * (total_attempted - 1)) + score) / total_attempted, 2)
        by_topic = quiz_stats.setdefault("by_topic", {})
        te = by_topic.setdefault(topic, {"attempts": 0, "avg_score": 0.0})
        te["attempts"] = int(te.get("attempts", 0)) + 1
        te["avg_score"] = round(((float(te.get("avg_score", 0)) * (te["attempts"] - 1)) + score) / te["attempts"], 2)
    if max_score > 0:
        _update_quiz_stats(stats, total_score, max_score)
    _update_streak(stats)
    _save_user_stats(stats, user_id)

    weak_topics_text: Optional[str] = None
    recommendation_text: Optional[str] = None
    if topic_scores:
        weak_topics_payload = {
            topic: round(sum(vals) / len(vals), 2) for topic, vals in topic_scores.items()
        }
        engine = get_ai_engine()
        try:
            weak_topic_response = engine.request(
                task_type=AITaskType.WEAK_TOPIC_DETECTION,
                user_input="Identify weak topics from this exam result.",
                context_data={
                    "exam_mode": "panic_mode",
                    "topic_scores": weak_topics_payload,
                    "total_questions": len(req.items),
                },
                offline_mode=True,
                privacy_sensitive=True,
                user_id=None,
            )
            weak_topics_text = weak_topic_response.text
            if weak_topics_text and weak_topics_text.strip():
                rec_response = engine.request(
                    task_type=AITaskType.STUDY_RECOMMENDATION,
                    user_input="Create a post-exam improvement plan.",
                    context_data={
                        "exam_mode": "panic_mode",
                        "topic_scores": weak_topics_payload,
                        "weak_topics_summary": weak_topics_text,
                        "study_time_minutes": 20,
                    },
                    offline_mode=True,
                    privacy_sensitive=True,
                    user_id=None,
                )
                recommendation_text = rec_response.text
        except (ConnectionError, TimeoutError):
            pass
        if not weak_topics_text:
            weak_topics_text = "AI unavailable for weak-topic analysis."
        if not recommendation_text:
            recommendation_text = "Complete more quizzes and review weak areas from your stats."

    # Enqueue for AWS sync (AppSync recordQuizAttempt) when sync enabled
    if req.results:
        _enqueue_panic_quiz_for_sync(req.results, len(req.items), user_id)

    return {"weak_topics_text": weak_topics_text, "recommendation_text": recommendation_text}


# ---------------------------------------------------------------------------
# User (auth-protected)
# ---------------------------------------------------------------------------


@app.get("/api/user/me")
def user_me(current_user: Annotated[User, Depends(get_current_user)]):
    """Return current authenticated user. Requires valid Bearer token."""
    return {
        "id": current_user.id,
        "username": current_user.username,
        "email": current_user.email,
    }


# ---------------------------------------------------------------------------
# User stats
# ---------------------------------------------------------------------------


@app.get("/api/user/stats")
def user_stats_get(current_user: Annotated[User, Depends(get_current_user)], user_id: str = Depends(get_user_id)):
    """Return user progress, streaks, preferences for the authenticated user.
    Reconciles flashcard_stats from actual deck data so Cards Mastered and Due reflect reality."""
    stats = _load_user_stats(user_id)
    ensure_streak_structure(stats)
    ensure_flashcard_structure(stats)
    decks = _load_flashcard_decks(user_id)
    cards = _all_cards_from_decks(decks)
    if cards:
        update_flashcard_stats_from_cards(stats, cards)
    _update_streak(stats)
    _save_user_stats(stats, user_id)
    return stats


def _build_insights_from_stats(stats: dict[str, Any], user_id: str) -> dict[str, Any]:
    """
    Build structured insights from user stats. Calls AI for weak-topic detection
    and study recommendation when data is available. Reuses logic from panic_finalize.
    """
    engine = get_ai_engine()
    streak_current = int(stats.get("streak", {}).get("current", 0) or 0)
    quiz_stats = stats.get("quiz_stats") or {}
    flashcard_stats = stats.get("flashcard_stats") or {}
    by_topic = quiz_stats.get("by_topic") or {}

    quiz_attempted = int(quiz_stats.get("total_attempted", 0) or 0)
    quiz_correct = int(quiz_stats.get("total_correct", 0) or 0)
    quiz_avg = float(quiz_stats.get("average_score", 0) or 0)
    quiz_accuracy = round((quiz_correct / quiz_attempted * 100)) if quiz_attempted > 0 else 0

    fc_reviewed = int(flashcard_stats.get("total_reviewed", 0) or 0)
    fc_mastered = int(flashcard_stats.get("mastered", 0) or 0)
    fc_mastery_pct = round((fc_mastered / fc_reviewed * 100)) if fc_reviewed > 0 else 0
    mastery_pct = round((quiz_accuracy * 0.6 + fc_mastery_pct * 0.4))

    topic_keys = list(by_topic.keys()) if isinstance(by_topic, dict) else []
    topic_scores_payload = (
        {t: round(float((by_topic.get(t) or {}).get("avg_score", 0) or 0), 2)
         for t in topic_keys}
        if topic_keys else {}
    )

    weak_topic_name = topic_keys[0] if topic_keys else "Not enough data"
    weak_topic_entry = by_topic.get(weak_topic_name) if isinstance(by_topic, dict) else None
    weak_score_raw = float((weak_topic_entry or {}).get("avg_score", 0) or 0)
    weak_topic_score = round(weak_score_raw * 10, 1) if weak_score_raw <= 1 else round(weak_score_raw * 10)

    weak_topics_text: Optional[str] = None
    study_recommendation_text: Optional[str] = None
    weak_topic_ai_name: Optional[str] = None

    if topic_scores_payload:
        try:
            weak_resp = engine.request(
                task_type=AITaskType.WEAK_TOPIC_DETECTION,
                user_input="Identify the single weakest topic from this student's quiz performance.",
                context_data={
                    "topic_scores": topic_scores_payload,
                    "total_questions": quiz_attempted,
                    "subject": "General",
                },
                offline_mode=True,
                privacy_sensitive=True,
                user_id=user_id,
            )
            weak_topics_text = weak_resp.text
            if weak_topics_text and weak_topics_text.strip():
                first_line = weak_topics_text.strip().split("\n")[0][:80]
                weak_topic_ai_name = first_line if first_line else weak_topic_name
        except (ConnectionError, TimeoutError):
            pass

        try:
            rec_resp = engine.request(
                task_type=AITaskType.STUDY_RECOMMENDATION,
                user_input="Create a personalized study plan for this student based on their stats.",
                context_data={
                    "topic_scores": topic_scores_payload,
                    "weak_topics_summary": weak_topics_text or "No weak topics identified.",
                    "study_time_minutes": 20,
                    "streak": streak_current,
                    "quiz_average": quiz_avg,
                    "flashcard_mastery_pct": fc_mastery_pct,
                    "total_quiz_attempted": quiz_attempted,
                },
                offline_mode=True,
                privacy_sensitive=True,
                user_id=user_id,
            )
            study_recommendation_text = rec_resp.text
        except (ConnectionError, TimeoutError):
            pass

    trend_points: list[float] = []
    quiz_history = _load_quiz_history(user_id)
    if quiz_history:
        percents = []
        for r in quiz_history[:10]:
            pct = r.get("percent")
            if pct is not None:
                percents.append(float(pct))
        trend_points = list(reversed(percents[-5:]))
    if len(trend_points) < 5:
        pad = [float(quiz_accuracy)] * (5 - len(trend_points))
        trend_points = trend_points + pad

    insights: list[dict[str, Any]] = [
        {
            "id": "insight_weak_topic",
            "title": "Weak Topic Alert",
            "description": f"Potential weak area: {weak_topic_ai_name or weak_topic_name}. Flag when below 60%.",
            "insight_type": "weak_topic_detection",
            "priority": "high",
            "related_subject": weak_topic_ai_name or weak_topic_name,
            "suggested_action": "Start a remedial quiz on this topic.",
            "weak_topic_name": weak_topic_ai_name or weak_topic_name,
            "weak_topic_score": weak_topic_score,
        },
        {
            "id": "insight_mastery",
            "title": "Subject Mastery Snapshot",
            "description": "Current mastery band: Beginner–Intermediate.",
            "insight_type": "subject_mastery",
            "priority": "medium",
            "related_subject": "All Subjects",
            "suggested_action": "Review weak concepts before the next quiz.",
            "mastery_pct": mastery_pct,
        },
        {
            "id": "insight_streak",
            "title": "Daily Learning Streak",
            "description": f"Current streak is {streak_current} day(s). Keep momentum going.",
            "insight_type": "daily_streak",
            "priority": "medium",
            "related_subject": "All Subjects",
            "suggested_action": "Complete one activity today to continue your streak.",
        },
        {
            "id": "insight_quiz_trend",
            "title": "Quiz Performance Trend",
            "description": f"Recent quiz accuracy is {quiz_accuracy}%.",
            "insight_type": "quiz_performance_trend",
            "priority": "medium",
            "related_subject": "Quiz",
            "suggested_action": "Retake a quiz in your weakest subject.",
            "trend_points": trend_points,
        },
        {
            "id": "insight_recommendation",
            "title": "AI Study Recommendation",
            "description": "Next best action based on weak topics and trend signals.",
            "insight_type": "study_recommendation",
            "priority": "low",
            "related_subject": "Recommended",
            "suggested_action": "Study 20 minutes + one focused quiz + flashcard recap.",
        },
    ]
    return {
        "insights": insights,
        "study_recommendation_text": study_recommendation_text,
    }


@app.get("/api/insights")
def insights_get(current_user: Annotated[User, Depends(get_current_user)], user_id: str = Depends(get_user_id)):
    """
    Return structured AI insights for the current user.
    Auth-protected. Uses real AI for weak-topic detection and study recommendation.
    """
    stats = _load_user_stats(user_id)
    ensure_streak_structure(stats)
    _update_streak(stats)
    result = _build_insights_from_stats(stats, user_id)
    return {
        "insights": result["insights"],
        "study_recommendation_text": result.get("study_recommendation_text"),
    }


@app.put("/api/user/stats")
def user_stats_put(stats: dict[str, Any], current_user: Annotated[User, Depends(get_current_user)], user_id: str = Depends(get_user_id)):
    """Update user progress/preferences for the authenticated user. Merges with existing."""
    existing = _load_user_stats(user_id)
    for key, value in stats.items():
        if isinstance(value, dict) and isinstance(existing.get(key), dict):
            existing[key] = {**existing[key], **value}
        else:
            existing[key] = value
    _save_user_stats(existing, user_id)
    return {"ok": True}


# ---------------------------------------------------------------------------
# Data export & clear (Settings: Export Data, Clear Local Data)
# ---------------------------------------------------------------------------


@app.get("/api/data/export")
def data_export(user_id: str = Depends(get_user_id)):
    """
    Export all user stats, flashcards, and profile for the authenticated user.
    """
    from datetime import datetime, timezone
    stats = _load_user_stats(user_id)
    cards = _load_flashcards(user_id)
    profile = load_profile_for_user(user_id)
    payload = {
        "exported_at": datetime.now(timezone.utc).isoformat(),
        "version": "1.0",
        "user_id": user_id,
        "user_stats": stats,
        "flashcards": cards,
        "profile": _profile_to_dict(profile),
    }
    return payload


@app.post("/api/data/clear")
def data_clear(user_id: str = Depends(get_user_id)):
    """
    Clear study data for the authenticated user: reset stats, clear flashcards, reset profile.
    Does NOT delete auth (users.db) or affect other users.
    """
    default = dict(_DEFAULT_STATS)
    default["user_id"] = user_id
    _save_user_stats(default, user_id)
    _save_flashcards([], user_id)
    existing = load_profile_for_user(user_id)
    reset_profile = UserProfile(
        profile_name=existing.profile_name if existing else user_id,
        profile_mode="solo",
        class_code=None,
        class_id=None,
        user_role="student",
        onboarding_complete=True,
    )
    save_profile_for_user(user_id, reset_profile)
    return {"ok": True, "message": f"Data cleared for user '{user_id}'."}


# ---------------------------------------------------------------------------
# User profile (AuthContext sync)
# ---------------------------------------------------------------------------


class ProfileRequest(BaseModel):
    """Profile update payload; all fields optional for partial merge."""
    profile_name: Optional[str] = None
    profile_mode: Optional[str] = None  # solo | teacher_linked | teacher_linked_provisional
    class_code: Optional[str] = None
    class_id: Optional[str] = None
    user_role: Optional[str] = None  # student | teacher
    onboarding_complete: Optional[bool] = None


def _profile_to_dict(p: Optional[UserProfile]) -> dict[str, Any]:
    """Convert UserProfile to JSON-serializable dict."""
    if p is None:
        return {"profile_name": None, "profile_mode": None, "class_code": None, "class_id": None, "user_role": None, "onboarding_complete": False}
    return {
        "profile_name": p.profile_name,
        "profile_mode": p.profile_mode,
        "class_code": p.class_code,
        "class_id": getattr(p, "class_id", None),
        "user_role": p.user_role,
        "onboarding_complete": getattr(p, "onboarding_complete", False),
    }


@app.get("/api/user/profile")
def user_profile_get(
    current_user: Annotated[User, Depends(get_current_user)],
    user_id: str = Depends(get_user_id),
):
    """Return persisted user profile scoped by JWT user_id. Auth-protected."""
    p = load_profile_for_user(user_id)
    return _profile_to_dict(p)


@app.post("/api/user/profile")
def user_profile_post(
    req: ProfileRequest,
    current_user: Annotated[User, Depends(get_current_user)],
    user_id: str = Depends(get_user_id),
):
    """Persist profile; merges with existing. Used by AuthContext. Returns saved profile."""
    existing = load_profile_for_user(user_id) or UserProfile()
    merged = UserProfile(
        profile_name=req.profile_name if req.profile_name is not None else existing.profile_name,
        profile_mode=req.profile_mode if req.profile_mode is not None else existing.profile_mode,
        class_code=req.class_code if req.class_code is not None else existing.class_code,
        class_id=req.class_id if req.class_id is not None else getattr(existing, "class_id", None),
        user_role=req.user_role if req.user_role is not None else existing.user_role,
        onboarding_complete=req.onboarding_complete if req.onboarding_complete is not None else getattr(existing, "onboarding_complete", False),
    )
    save_profile_for_user(user_id, merged)
    return _profile_to_dict(merged)


# ---------------------------------------------------------------------------
# Sync (SyncManager + ConflictAwareOrchestrator)
# ---------------------------------------------------------------------------



def _get_conflict_engine(user_id: str):
    """Create a user-scoped ConflictResolutionEngine for conflict endpoints."""
    try:
        from conflict_resolution_engine import ConflictResolutionEngine
        return ConflictResolutionEngine(base_path=str(BASE_PATH), user_id=user_id)
    except ImportError:
        return None


def _persist_resolved_entity(entity_type: str, entity_id: str, resolved_data: dict, user_id: str) -> None:
    """Persist resolved conflict data to local store based on entity type."""
    stats = _load_user_stats(user_id)
    et = (entity_type or "").lower()
    if et in ("userstats", "user_stats") or entity_id in ("user_stats", "stats"):
        _save_user_stats(resolved_data, user_id)
    elif et in ("streakrecord", "streak"):
        merged = dict(stats)
        merged["streak"] = {**(stats.get("streak") or {}), **resolved_data}
        _save_user_stats(merged, user_id)
    elif et in ("quizstats", "quiz_stats"):
        merged = dict(stats)
        merged["quiz_stats"] = {**(stats.get("quiz_stats") or {}), **resolved_data}
        _save_user_stats(merged, user_id)
    else:
        # Default: deep merge top-level keys into user_stats
        for k, v in resolved_data.items():
            if k in stats and isinstance(stats[k], dict) and isinstance(v, dict):
                stats[k] = {**stats[k], **v}
            else:
                stats[k] = v
        _save_user_stats(stats, user_id)


@app.post("/api/sync")
def sync_trigger(current_user: Annotated[User, Depends(get_current_user)], user_id: str = Depends(get_user_id)):
    """Trigger sync with AWS when online. Uses SyncManager.try_sync()."""
    try:
        from sync_manager import SyncManager
        sm = SyncManager(base_path=str(BASE_PATH), user_id=user_id)
        result = sm.try_sync()
        return {
            "ok": True,
            "synced": result.get("synced", 0),
            "failed": result.get("failed", 0),
            "pending": result.get("pending", 0),
            "online": result.get("online", False),
            "aws_sync": result.get("aws_sync"),
            "errors": result.get("errors", []),
            "message": (
                f"Synced {result.get('synced', 0)} item(s)"
                if result.get("synced", 0) > 0
                else ("Sync complete" if not result.get("pending") else "Some items pending")
            ),
        }
    except ImportError as e:
        return {"ok": False, "message": f"SyncManager unavailable: {e}", "synced": 0, "failed": 0, "pending": 0, "online": False, "errors": []}
    except Exception as e:
        return {"ok": False, "message": str(e), "synced": 0, "failed": 0, "pending": 0, "online": False, "errors": [str(e)]}


@app.get("/api/sync/status")
def sync_status(current_user: Annotated[User, Depends(get_current_user)], user_id: str = Depends(get_user_id)):
    """Return sync status: queue summary, connectivity, last sync."""
    stats = _load_user_stats(user_id)
    prefs = stats.get("preferences") or {}
    sync_enabled = prefs.get("sync_enabled", True)
    last_sync = stats.get("last_sync_timestamp")

    out: dict[str, Any] = {
        "sync_enabled": sync_enabled,
        "last_sync_timestamp": last_sync,
        "online": False,
        "queue": {"total": 0, "quiz_attempts": 0, "streak_updates": 0, "oldest_item": None},
    }

    if not sync_enabled:
        return out

    try:
        from sync_manager import SyncManager
        sm = SyncManager(base_path=str(BASE_PATH), user_id=user_id)
        out["online"] = sm.check_connectivity()
        out["queue"] = sm.get_queue_summary()
    except Exception:
        pass

    return out


@app.get("/api/sync/conflicts")
def get_conflicts(current_user: Annotated[User, Depends(get_current_user)], user_id: str = Depends(get_user_id)):
    """Return pending sync conflicts for the authenticated user (scoped by user_id from JWT)."""
    engine = _get_conflict_engine(user_id)
    if engine is None:
        return {"conflicts": [], "message": "Conflict resolution engine unavailable"}
    conflicts = engine.get_pending_conflicts()
    # Ensure reason is string for JSON
    for c in conflicts:
        if hasattr(c.get("reason"), "value"):
            c["reason"] = c["reason"].value
    return {"conflicts": conflicts}


class ResolveConflictRequest(BaseModel):
    choice: str = Field(..., description="keep_local | keep_cloud | merge")


@app.post("/api/sync/conflicts/{entity_id}/resolve")
def resolve_conflict(entity_id: str, body: ResolveConflictRequest, current_user: Annotated[User, Depends(get_current_user)], user_id: str = Depends(get_user_id)):
    """Resolve a conflict by entity_id. Persists resolved data and removes from pending."""
    choice = (body.choice or "").strip().lower()
    if choice not in ("keep_local", "keep_cloud", "merge"):
        raise HTTPException(400, "choice must be keep_local, keep_cloud, or merge")

    engine = _get_conflict_engine(user_id)
    if engine is None:
        raise HTTPException(503, "Conflict resolution engine unavailable")

    try:
        from conflict_resolution_engine import ConflictResult
        conflicts = engine.get_pending_conflicts()
        conflict_dict = next((c for c in conflicts if c.get("entity_id") == entity_id), None)
        entity_type = conflict_dict.get("entity_type", "UserStats") if conflict_dict else "UserStats"

        if not conflict_dict:
            raise ValueError(f"No pending conflict found for entity {entity_id}")
        conflict = ConflictResult.from_dict(conflict_dict)
        resolved_data = engine.apply_manual_resolution(conflict, choice)
        _persist_resolved_entity(entity_type, entity_id, resolved_data, user_id)
        try:
            from sync_manager import SyncManager
            SyncManager(base_path=str(BASE_PATH), user_id=user_id).try_sync()
        except Exception:
            pass
        return {"ok": True, "entity_id": entity_id, "choice": choice}
    except ValueError as e:
        raise HTTPException(404, str(e))
    except Exception as e:
        raise HTTPException(500, str(e))


# ---------------------------------------------------------------------------
# RAG search (ChromaDB semantic search)
# ---------------------------------------------------------------------------

_rag_vector_store: Optional[Any] = None


def _get_rag_vector_store():
    """Lazy-load ChromaDB vector store. Returns None if DB unavailable."""
    global _rag_vector_store
    if _rag_vector_store is not None:
        return _rag_vector_store
    try:
        from ai_chat.vector import build_vector_store
        _rag_vector_store = build_vector_store(rebuild=False)
        return _rag_vector_store
    except ImportError as e:
        return None
    except Exception:
        return None


def _get_relevant_chunks_from_chromadb(
    query: str, k: int = 5, source_filter: Optional[str] = None, max_chars: int = 12_000
) -> str:
    """
    Run semantic search over ChromaDB and return concatenated chunk text for use as source_content.
    If source_filter is set, only chunks from that textbook (source metadata) are returned.
    Returns empty string if store unavailable or on error (caller can fall back to file-based flow).
    """
    if not query or not query.strip():
        return ""
    store = _get_rag_vector_store()
    if store is None:
        return ""
    search_kwargs: dict[str, Any] = {"k": k}
    if source_filter:
        search_kwargs["filter"] = {"source": source_filter}
    try:
        retriever = store.as_retriever(search_kwargs=search_kwargs)
        docs = retriever.invoke(query.strip())
    except Exception:
        return ""
    parts: list[str] = []
    total = 0
    for doc in docs or []:
        content = (getattr(doc, "page_content", None) or str(doc)).strip()
        if not content:
            continue
        if total + len(content) + 2 > max_chars:
            parts.append(content[: max_chars - total - 2].rstrip())
            break
        parts.append(content)
        total += len(content) + 2
    return "\n\n".join(parts) if parts else ""


def _rag_search(query: str, k: int = 5) -> tuple[list[dict[str, Any]], Optional[str]]:
    """
    Run semantic search over ChromaDB. Returns (results, error_message).
    On success, error_message is None. On failure, results is [] and error_message describes the issue.
    """
    if not query or not query.strip():
        return [], "Provide a non-empty query."

    store = _get_rag_vector_store()
    if store is None:
        return [], "ChromaDB unavailable. Install langchain-chroma and langchain-huggingface, and ensure the vector store is initialized."

    try:
        retriever = store.as_retriever(search_kwargs={"k": k})
        docs = retriever.invoke(query.strip())
    except Exception as e:
        return [], f"ChromaDB query failed: {str(e)}"

    results: list[dict[str, Any]] = []
    for doc in docs or []:
        content = getattr(doc, "page_content", None) or str(doc)
        metadata = getattr(doc, "metadata", None) or {}
        if isinstance(metadata, dict):
            results.append({
                "content": content[:2000] + ("..." if len(content) > 2000 else ""),
                "source": metadata.get("source", "unknown"),
                "subject": metadata.get("subject", "general"),
            })
        else:
            results.append({"content": content[:2000], "source": "unknown", "subject": "general"})

    return results, None


@app.get("/api/rag/search")
def rag_search(q: str = "", k: int = 5):
    """Semantic search over local ChromaDB. Returns top-k matching chunks from embedded textbooks."""
    results, err = _rag_search(q, k=min(max(1, k), 20))
    if err:
        return {"results": [], "message": err}
    return {"results": results}


# ---------------------------------------------------------------------------
# Serve compiled React SPA from frontend/dist (when running as backend.main:app)
# Repo root main.py also mounts SPA when using python main.py.
# IMPORTANT: Do NOT mount StaticFiles at "/" — it catches all requests (including
# /api/*) and returns 405 for POST. Serve /assets and use a GET catch-all for SPA.
# ---------------------------------------------------------------------------

_DIST = _APP_DIR.parent / "frontend" / "dist"

if _DIST.is_dir():
    # Static assets (Vite puts JS/CSS in /assets/)
    _ASSETS = _DIST / "assets"
    if _ASSETS.is_dir():
        app.mount("/assets", StaticFiles(directory=str(_ASSETS)), name="spa_assets")

    # Favicon at root
    _VITE_SVG = _DIST / "vite.svg"
    if _VITE_SVG.is_file():

        @app.get("/vite.svg")
        def _serve_favicon():
            return FileResponse(_VITE_SVG)

    # SPA fallback: GET catch-all returns index.html (client-side routing)
    # API routes are defined above, so POST /api/* etc. match before this
    @app.get("/{full_path:path}")
    def _serve_spa(full_path: str):
        return FileResponse(_DIST / "index.html")
